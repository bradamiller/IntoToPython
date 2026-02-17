#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 5: Proportional Control with One Sensor."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

# ── Fonts ──
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

# ── Sizes ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              row_height=Inches(0.45)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        cw = int(width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = cw

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5),
        Inches(12), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(3.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 5: Proportional Control with One Sensor"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    objectives = [
        "Understand what a setpoint is and how to calculate error",
        "Use set_effort() for continuous motor control",
        "Build a proportional control loop for line following",
        "Tune the Kp gain constant by observing behavior",
    ]
    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5.5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    agenda = [
        ("The problem with bounce driving", "5 min"),
        ("set_effort() and continuous control", "5 min"),
        ("Proportional controller step by step", "20 min"),
        ("Follow the circle edge, tune Kp", "20 min"),
    ]
    for item, duration in agenda:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(16), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: Hook — What is Wrong with Bounce Driving?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "What is Wrong with Bounce Driving?")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Bounce driving (Lessons 3\u20134):", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Drive straight \u2192 Hit line \u2192 Stop \u2192 Turn \u2192 Drive straight \u2192 Hit line \u2192 \u2026",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Problems:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Jerky, not smooth", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Robot is NOT following the line \u2014 it is bouncing off it", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Cannot handle curves well", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "What we want:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot smoothly follows the line", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Constant small adjustments, like steering a car", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, 'Question: "How can the robot steer while it is driving?"',
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 3: Blocking vs. Continuous Motor Control
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Blocking vs. Continuous Motor Control")

    add_label(slide, "Blocking calls (what we used before):",
              Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.straight(30)   # Drives 30 cm, then STOPS\n"
        "drivetrain.turn(90)       # Turns 90 degrees, then STOPS",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.4)
    )
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(6.0), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_text(tf, "Robot does one thing at a time", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Cannot read sensors WHILE driving", level=0, font_size=Pt(16), color=DARK_GRAY)

    add_label(slide, "Continuous control (new):",
              Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.set_effort(0.3, 0.3)   # Sets motors and KEEPS GOING",
        Inches(7.0), Inches(2.0), Inches(6.0), Inches(0.9)
    )
    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(3.0), Inches(6.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_text(tf, "Sets motor power and returns immediately", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Motors keep running until you change them", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "You can read sensors in between calls", level=0, font_size=Pt(16), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key difference: set_effort() does not block \u2014 the program keeps running."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 4: set_effort() Examples
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "set_effort() Examples")

    add_label(slide, "Drive straight:", Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.set_effort(0.3, 0.3)   # Both motors equal",
        Inches(0.6), Inches(2.0), Inches(12.0), Inches(0.9)
    )

    add_label(slide, "Turn right (left faster):", Inches(0.6), Inches(3.0), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.set_effort(0.4, 0.2)   # Left faster than right",
        Inches(0.6), Inches(3.6), Inches(12.0), Inches(0.9)
    )

    add_label(slide, "Turn left (right faster):", Inches(0.6), Inches(4.6), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.set_effort(0.2, 0.4)   # Right faster than left",
        Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.9)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Values range from -1.0 (full reverse) to 1.0 (full forward)    |    set_effort(0.5, -0.5) = spin in place"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 5: The Setpoint
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Setpoint")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "From Lesson 1, you recorded sensor values:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "On white surface: ~0.1 to 0.3", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "On the edge of the line: ~0.5", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "On the black line: ~0.7 to 0.9", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "The setpoint is the target value we want to maintain.",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)

    add_code_block(slide,
        "setpoint = 0.5   # The edge of the line",
        Inches(0.6), Inches(4.3), Inches(6.0), Inches(0.9)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Goal: Keep the sensor reading at the setpoint by making constant small steering adjustments."
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 6: Calculating Error
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Calculating Error")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Error = how far we are from where we want to be"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "error = setpoint - sensor_reading",
        Inches(0.6), Inches(2.0), Inches(7.0), Inches(0.9)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Examples (setpoint = 0.5):"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    headers = ["Where is the robot?", "sensor_reading", "error", "Meaning"]
    rows = [
        ["On the edge (perfect)", "0.5", "0.0", "No correction needed"],
        ["Drifted onto white",    "0.2", "+0.3", "Steer toward line"],
        ["Drifted onto black",    "0.8", "\u22120.3", "Steer away from line"],
    ]
    col_widths = [Inches(3.2), Inches(2.2), Inches(1.5), Inches(4.0)]
    add_table(slide, headers, rows, Inches(0.5), Inches(3.7), Inches(11.0),
              col_widths=col_widths, row_height=Inches(0.42))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Positive error = robot is too far off the line (on white) \u2014 steer toward it",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Negative error = robot is too far on the line (on black) \u2014 steer away",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Zero error = perfect position \u2014 drive straight",
                    level=0, font_size=Pt(16), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 7: From Error to Correction
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Error to Correction")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "We multiply error by a gain constant Kp:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "correction = error * Kp",
        Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.9)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Kp controls how aggressively the robot reacts.    Example with Kp = 0.5:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    headers = ["error", "correction = error \u00d7 0.5", "What happens"]
    rows = [
        ["+0.3", "+0.15", "Gentle steer toward line"],
        ["\u22120.3", "\u22120.15", "Gentle steer away from line"],
        ["+0.4", "+0.20", "Stronger steer toward line"],
        ["0.0", "0.0", "Drive straight"],
    ]
    col_widths = [Inches(1.8), Inches(3.8), Inches(5.0)]
    add_table(slide, headers, rows, Inches(0.5), Inches(3.6), Inches(10.5),
              col_widths=col_widths, row_height=Inches(0.42))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Then apply correction to motors:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "left_effort  = base_effort + correction\n"
        "right_effort = base_effort - correction",
        Inches(0.6), Inches(6.4), Inches(7.5), Inches(1.0)
    )

    # ════════════════════════════════════════
    # SLIDE 8: The Complete Control Loop
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Complete Control Loop")

    add_code_block(slide,
        "setpoint = 0.5\n"
        "Kp = 0.5\n"
        "base_effort = 0.3\n"
        "\n"
        "while True:\n"
        "    # 1. Read sensor\n"
        "    sensor_value = reflectance.get_left()\n"
        "\n"
        "    # 2. Calculate error\n"
        "    error = setpoint - sensor_value\n"
        "\n"
        "    # 3. Calculate correction\n"
        "    correction = error * Kp\n"
        "\n"
        "    # 4. Apply to motors\n"
        "    left_effort = base_effort + correction\n"
        "    right_effort = base_effort - correction\n"
        "    drivetrain.set_effort(left_effort, right_effort)\n"
        "\n"
        "    time.sleep(0.01)",
        Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.8),
        font_size=Pt(13)
    )

    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(1.4), Inches(4.8), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "The loop:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Read Sensor", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "\u2193", level=0, font_size=Pt(18), color=MEDIUM_GRAY)
    add_bullet_text(tf, "Calculate Error", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "\u2193", level=0, font_size=Pt(18), color=MEDIUM_GRAY)
    add_bullet_text(tf, "Calculate Correction", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "\u2193", level=0, font_size=Pt(18), color=MEDIUM_GRAY)
    add_bullet_text(tf, "Set Motors", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "\u2193 (back to top)", level=0, font_size=Pt(18), color=MEDIUM_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Runs hundreds of times per second!",
                    level=0, font_size=Pt(17), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 9: Tracing Through the Math
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Tracing Through the Math")

    add_label(slide, "Scenario: Robot drifts onto white  (sensor reads 0.2)",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
              font_size=Pt(18), color=DARK_BLUE)

    add_code_block(slide,
        "setpoint      = 0.5\n"
        "sensor_value  = 0.2\n"
        "error         = 0.5 - 0.2 = 0.3\n"
        "Kp            = 0.5\n"
        "correction    = 0.3 * 0.5 = 0.15\n"
        "base_effort   = 0.3\n"
        "\n"
        "left_effort   = 0.3 + 0.15 = 0.45\n"
        "right_effort  = 0.3 - 0.15 = 0.15",
        Inches(0.6), Inches(1.9), Inches(6.5), Inches(3.5)
    )

    txBox = slide.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Result: Left motor (0.45) faster than right (0.15) \u2192 robot steers right, back toward the line."
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    p.font.name = BODY_FONT

    add_label(slide, "Scenario: Robot drifts onto black  (sensor reads 0.8)",
              Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
              font_size=Pt(18), color=DARK_BLUE)

    add_code_block(slide,
        "error         = 0.5 - 0.8 = -0.3\n"
        "correction    = -0.3 * 0.5 = -0.15\n"
        "\n"
        "left_effort   = 0.3 + (-0.15) = 0.15\n"
        "right_effort  = 0.3 - (-0.15) = 0.45",
        Inches(0.6), Inches(6.2), Inches(6.5), Inches(2.2)
    )

    txBox = slide.shapes.add_textbox(Inches(7.3), Inches(6.2), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Result: Right motor (0.45) faster \u2192 robot steers left, away from the line."
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 10: What Does Kp Do?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "What Does Kp Do?")

    headers = ["Kp Value", "Behavior"]
    rows = [
        ["Too low (e.g., 0.1)", "Robot reacts too slowly, drifts off the line"],
        ["Just right (e.g., 0.5)", "Robot follows smoothly with small adjustments"],
        ["Too high (e.g., 2.0)", "Robot overcorrects, oscillates wildly"],
    ]
    col_widths = [Inches(3.5), Inches(7.5)]
    add_table(slide, headers, rows, Inches(0.6), Inches(1.5), Inches(11.0),
              col_widths=col_widths, row_height=Inches(0.55))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Tuning Kp is about finding the sweet spot:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Start at 0.5", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "If too wiggly, lower it", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "If too sluggish, raise it", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 11: Exercise — Follow the Circle
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Exercise \u2014 Follow the Circle")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.2), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Your task:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "1. Type in the complete program", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "2. Place the robot on the edge of the taped circle", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "3. Press the button and observe", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "4. Tune Kp and base_effort for smooth tracking", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Success:", level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)
    add_bullet_text(tf, "Robot completes one full lap of the circle.", level=0, font_size=Pt(17), color=DARK_GRAY)

    add_code_block(slide,
        "from XRPLib.differential_drive import DifferentialDrive\n"
        "from XRPLib.reflectance import Reflectance\n"
        "from XRPLib.board import Board\n"
        "import time\n"
        "\n"
        "drivetrain  = DifferentialDrive.get_default_differential_drive()\n"
        "reflectance = Reflectance.get_default_reflectance()\n"
        "board       = Board.get_default_board()\n"
        "\n"
        "setpoint    = 0.5\n"
        "Kp          = 0.5\n"
        "base_effort = 0.3\n"
        "\n"
        "board.wait_for_button()\n"
        "\n"
        "while True:\n"
        "    sensor_value = reflectance.get_left()\n"
        "    error        = setpoint - sensor_value\n"
        "    correction   = error * Kp\n"
        "    drivetrain.set_effort(base_effort + correction,\n"
        "                          base_effort - correction)\n"
        "    time.sleep(0.01)",
        Inches(6.0), Inches(1.4), Inches(7.0), Inches(5.8),
        font_size=Pt(12)
    )

    # ════════════════════════════════════════
    # SLIDE 12: Connection to Lesson 6
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Lesson 6")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "What we accomplished today:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Smooth line following with ONE sensor using proportional control",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Understood setpoint, error, Kp, and correction",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Limitation of one sensor:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Only tracks one edge of the line", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "If the robot drifts too far off, it cannot tell which side it is on",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next lesson: Two-Sensor Line Following", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Use BOTH left and right sensors", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Error = left sensor \u2212 right sensor", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Even smoother tracking, follows the CENTER of the line",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, 'Question: "Why might two sensors be better than one?"',
                    level=0, font_size=Pt(19), bold=True, color=ACCENT_BLUE)

    add_code_block(slide,
        "error = reflectance.get_left() - reflectance.get_right()",
        Inches(0.6), Inches(6.5), Inches(8.5), Inches(0.8)
    )

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "05-proportional-control.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
