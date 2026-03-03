"""Build the Proportional Control Review slide deck.

Uses the existing 05-proportional-control.pptx as a template for consistent
styling (header bars, fonts, code blocks, tables).
Run: /opt/homebrew/anaconda3/bin/python3 build_review_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Constants matching existing deck styling ─────────────────────────────────
SLIDE_W = 12191695  # 13.33 in (widescreen 16:9)
SLIDE_H = 6858000   # 7.50 in

HEADER_COLOR = RGBColor(0x1B, 0x3A, 0x5C)  # dark navy
HEADER_H = 685800   # ~0.75 in (shorter than title slide)
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)
CODE_BORDER = RGBColor(0x99, 0x99, 0x99)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)

MARGIN_LEFT = Inches(0.6)
CONTENT_W = Inches(12.0)

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"
FONT_TITLE = "Calibri"


# ── Helper functions ─────────────────────────────────────────────────────────

def add_header_bar(slide, height=HEADER_H):
    """Add the dark navy header bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_COLOR
    shape.line.fill.background()
    return shape


def add_title(slide, text, top=Inches(0.12), size=Pt(32)):
    """Add white title text over the header bar."""
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_W, Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TITLE
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = WHITE
    return txbox


def add_text(slide, text, left=MARGIN_LEFT, top=Inches(1.0), width=CONTENT_W,
             height=Inches(0.5), size=Pt(18), bold=False, color=DARK_GRAY,
             alignment=None, font_name=FONT_BODY):
    """Add a simple text box."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txbox


def add_bullet_list(slide, items, left=MARGIN_LEFT, top=Inches(1.5),
                    width=CONTENT_W, size=Pt(16), color=DARK_GRAY,
                    bold_items=None):
    """Add a bulleted text list. bold_items is a set of indices to bold."""
    height = Inches(0.4 * len(items))
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    bold_items = bold_items or set()

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"  \u2022  {item}"
        run.font.name = FONT_BODY
        run.font.size = size
        run.font.bold = i in bold_items
        run.font.color.rgb = color
    return txbox


def add_code_block(slide, code, left=MARGIN_LEFT, top=Inches(2.0),
                   width=Inches(11.0), size=Pt(14)):
    """Add code in a rounded-rectangle code block."""
    lines = code.strip().split("\n")
    line_h = size.pt * 1.4
    height = Emu(int((len(lines) * line_h + 20) * 12700))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CODE_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_CODE
        run.font.size = size
        run.font.color.rgb = DARK_GRAY
    return shape


def add_table(slide, headers, rows, left=MARGIN_LEFT, top=Inches(2.5),
              width=Inches(11.5), row_height=Inches(0.35), font_size=Pt(14),
              header_color=ACCENT_BLUE, alt_row_color=LIGHT_BLUE):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    col_w = int(width / num_cols)

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = tbl_shape.table

    # Set column widths evenly
    for ci in range(num_cols):
        table.columns[ci].width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT_BODY
                run.font.size = font_size
                run.font.bold = True
                run.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_color
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = font_size
                    run.font.color.rgb = DARK_GRAY

    return tbl_shape


# ── Build the presentation ───────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank


# ━━━ SLIDE 1: Title ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
# Large header bar for title slide
add_header_bar(sl, height=Emu(2743200))
add_text(sl, "Module 2: Line Tracking", top=Inches(0.5), size=Pt(22),
         color=WHITE, left=MARGIN_LEFT)
add_text(sl, "Proportional Control Review", top=Inches(1.2), size=Pt(36),
         bold=True, color=WHITE, left=MARGIN_LEFT)
add_text(sl, "Before we move on, let's make sure we REALLY understand\n"
             "what proportional control is doing \u2014 not just the code, but the idea.",
         top=Inches(3.2), size=Pt(20), color=DARK_GRAY)


# ━━━ SLIDE 2: Shower Analogy ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "The Shower Analogy")
add_text(sl, "Proportional control is something you already do every day.",
         top=Inches(0.9), size=Pt(20), bold=True)

add_table(sl,
    ["Situation", "How you feel", "What you do"],
    [
        ["Water is WAY too cold", "Very uncomfortable", "Turn knob a LOT toward hot"],
        ["Water is a little too cold", "Slightly uncomfortable", "Turn knob a LITTLE toward hot"],
        ["Water is just right", "Comfortable", "Do nothing"],
        ["Water is a little too hot", "Slightly uncomfortable", "Turn knob a LITTLE toward cold"],
        ["Water is WAY too hot", "Very uncomfortable", "Turn knob a LOT toward cold"],
    ],
    top=Inches(1.5), row_height=Inches(0.4))

add_text(sl, "The key idea:  How much you turn the knob depends on how far off the temperature is.\n"
             "A little off?  Small adjustment.     Way off?  Big adjustment.\n"
             "That is proportional control.  The response is PROPORTIONAL to the error.",
         top=Inches(4.5), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 3: Same Idea for the Robot ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "The Same Idea, Applied to the Robot")

add_code_block(sl,
    "  White surface              LINE              White surface\n"
    "        A         B         C         D         E\n"
    "     (far       (a bit    (on the   (a bit    (far\n"
    "     left)       left)     edge)    right)    right)",
    top=Inches(1.0), size=Pt(16))

add_table(sl,
    ["Position", "Sensor reads", "How far off?", "Robot does"],
    [
        ["A  (far left of line)", "0.1", "Very far", "Steer HARD toward line"],
        ["B  (a little left)", "0.3", "A little", "Steer gently toward line"],
        ["C  (on the edge)", "0.5", "Perfect", "Drive straight"],
        ["D  (a little right)", "0.7", "A little", "Steer gently away from line"],
        ["E  (far right of line)", "0.9", "Very far", "Steer HARD away from line"],
    ],
    top=Inches(3.2), row_height=Inches(0.38))

add_text(sl, "Same principle as the shower:  bigger error = bigger correction",
         top=Inches(6.0), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 4: What is Error? ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, 'What is "Error"?')
add_text(sl, "Error = how far am I from where I want to be?",
         top=Inches(0.9), size=Pt(20), bold=True)

add_code_block(sl, "error = setpoint - sensor_reading", top=Inches(1.5), size=Pt(20),
               width=Inches(8))

add_code_block(sl,
    " 0.0       0.2       0.4    0.5    0.6       0.8       1.0\n"
    "  |----white----|---------|---*---|---------|----black----|\n"
    "                           setpoint\n"
    "\n"
    "  Sensor reads 0.2:  error = 0.5 - 0.2 = +0.3   (steer toward line -->)\n"
    "  Sensor reads 0.5:  error = 0.5 - 0.5 =  0.0   (perfect, go straight)\n"
    "  Sensor reads 0.8:  error = 0.5 - 0.8 = -0.3   (<-- steer away from line)",
    top=Inches(2.5), size=Pt(15))

add_bullet_list(sl, [
    "The SIGN  (+/-)  tells you WHICH DIRECTION to steer",
    "The SIZE  (big/small)  tells you HOW MUCH to steer",
], top=Inches(5.5), size=Pt(18), bold_items={0, 1})


# ━━━ SLIDE 5: What is Kp? ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "What is Kp?")
add_text(sl, "Kp controls how aggressively the robot reacts.  Think of Kp as the robot's personality:",
         top=Inches(0.9), size=Pt(18))

add_table(sl,
    ["Kp value", "Personality", "Behavior"],
    [
        ["Low  (0.1)", '\U0001f634  "Lazy robot"', "Barely reacts, drifts off the line"],
        ["Medium  (0.5)", '\U0001f642  "Calm robot"', "Smooth, steady corrections"],
        ["High  (2.0)", '\U0001f630  "Nervous robot"', "Overreacts, zigzags wildly"],
    ],
    top=Inches(1.5), row_height=Inches(0.45), font_size=Pt(16))

add_code_block(sl, "correction = error * Kp", top=Inches(3.3), size=Pt(20), width=Inches(6))

add_text(sl, "Same error, different Kp:", top=Inches(4.2), size=Pt(16), bold=True)
add_table(sl,
    ["error", "Kp = 0.1", "Kp = 0.5", "Kp = 2.0"],
    [
        ["+0.3", "0.03  (tiny steer)", "0.15  (moderate)", "0.60  (HUGE steer)"],
        ["-0.3", "-0.03  (tiny steer)", "-0.15  (moderate)", "-0.60  (HUGE steer)"],
    ],
    top=Inches(4.7), row_height=Inches(0.38))

add_text(sl, "Kp is a volume knob for the correction.",
         top=Inches(5.9), size=Pt(20), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 6: What Does This Look Like on the Track? ━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "What Does This Look Like on the Track?")

# Three labeled path diagrams as code blocks with labels
y = Inches(1.1)
for label, path, desc, bg in [
    ('Kp too low (0.1)  \u2014  "Lazy robot"',
     "Line:   ~~~~~~~~~~~~~~~~~~~~~~~~~~~~\n"
     "Robot:  ~~~~~~~~~~~~~~~\n"
     "                       ~~~~~~~  (drifts off)",
     "Robot reacts too slowly, drifts off on curves", LIGHT_RED),
    ('Kp just right (0.5)  \u2014  "Calm robot"',
     "Line:   ~~~~~~~~~~~~~~~~~~~~~~~~~~~~\n"
     "Robot:  ~~~~~~~~~~~~~~~~~~~~~~~~~~~~  (smooth!)",
     "Robot follows smoothly with small adjustments", LIGHT_GREEN),
    ('Kp too high (2.0)  \u2014  "Nervous robot"',
     "Line:   ~~~~~~~~~~~~~~~~~~~~~~~~~~~~\n"
     "Robot:  ~/\\/\\/\\/\\/\\/\\/\\/\\/\\/\\/\\/\\~  (zigzag!)",
     "Robot overcorrects, oscillates wildly", LIGHT_YELLOW),
]:
    add_text(sl, label, top=y, size=Pt(16), bold=True, color=ACCENT_BLUE)
    y += Inches(0.35)
    block = add_code_block(sl, path, top=y, size=Pt(13), width=Inches(8))
    add_text(sl, desc, top=y + Inches(0.05), left=Inches(9.0), width=Inches(3.5),
             size=Pt(14), color=MED_GRAY)
    y += Inches(1.3)

add_text(sl, "The goal of tuning: find the Kp where the robot follows smoothly without oscillating.",
         top=Inches(6.1), size=Pt(16), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 7: The Four Steps ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "The Four Steps \u2014 Every Single Loop")

add_text(sl, "The control loop does the same four steps, hundreds of times per second:",
         top=Inches(0.9), size=Pt(18))

add_code_block(sl,
    "  +---> 1. READ SENSOR ----> 2. CALCULATE ERROR ---+\n"
    "  |       (where am I?)        (how far off?)       |\n"
    "  |                                                  |\n"
    "  +--- 4. SET MOTORS <---- 3. CALC CORRECTION ------+\n"
    "       (adjust steering)     (how much to fix?)",
    top=Inches(1.6), size=Pt(16))

add_text(sl, "Trace it with real numbers:", top=Inches(3.5), size=Pt(18), bold=True)

add_code_block(sl,
    "1. Read sensor:          sensor_value = 0.2\n"
    "2. Calculate error:      error = 0.5 - 0.2 = +0.3\n"
    "3. Calculate correction: correction = 0.3 * 0.5 = 0.15\n"
    "4. Set motors:           arcade(0.3, 0.15)\n"
    "                         --> left motor  = 0.3 + 0.15 = 0.45\n"
    "                         --> right motor = 0.3 - 0.15 = 0.15\n"
    "                         --> robot steers right, toward the line",
    top=Inches(4.0), size=Pt(15))

add_text(sl, "Then it does it again.  And again.  And again.  Hundreds of times per second.",
         top=Inches(6.2), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 8: Let's Trace It Together (blank table) ━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Let's Trace It Together")

add_text(sl, "Settings:  setpoint = 0.5     Kp = 0.5     base_effort = 0.3",
         top=Inches(0.9), size=Pt(16), font_name=FONT_CODE)

add_table(sl,
    ["sensor", "error = 0.5 - sensor", "correction = error * 0.5",
     "Left motor", "Right motor", "Which way?"],
    [
        ["0.5", "?", "?", "?", "?", "?"],
        ["0.2", "?", "?", "?", "?", "?"],
        ["0.8", "?", "?", "?", "?", "?"],
        ["0.1", "?", "?", "?", "?", "?"],
        ["0.6", "?", "?", "?", "?", "?"],
    ],
    top=Inches(1.5), row_height=Inches(0.4), font_size=Pt(15))

add_text(sl, "Work through each row \u2014 calculate error, correction, and motor values.",
         top=Inches(4.2), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 9: Answers Revealed ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Answers Revealed")

add_text(sl, "Settings:  setpoint = 0.5     Kp = 0.5     base_effort = 0.3",
         top=Inches(0.9), size=Pt(16), font_name=FONT_CODE)

add_table(sl,
    ["sensor", "error", "correction", "Left motor", "Right motor", "Which way?"],
    [
        ["0.5", " 0.0", " 0.0", "0.30", "0.30", "Straight"],
        ["0.2", "+0.3", "+0.15", "0.45", "0.15", "Steer right (toward line)"],
        ["0.8", "-0.3", "-0.15", "0.15", "0.45", "Steer left (away from line)"],
        ["0.1", "+0.4", "+0.20", "0.50", "0.10", "Steer hard right"],
        ["0.6", "-0.1", "-0.05", "0.25", "0.35", "Gentle steer left"],
    ],
    top=Inches(1.5), row_height=Inches(0.4), font_size=Pt(15))

add_text(sl, "Notice:  bigger error \u2192 bigger correction \u2192 bigger difference between motors",
         top=Inches(4.2), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 10: Two Sensors ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Two Sensors \u2014 Even Simpler")

add_code_block(sl, "error = left_sensor - right_sensor     # No setpoint needed!",
               top=Inches(1.0), size=Pt(20), width=Inches(10))

add_code_block(sl,
    "Centered:              Drifted left:           Drifted right:\n"
    "  [L]  LINE  [R]         [L]  LINE       [R]     [L]       LINE  [R]\n"
    "  0.5        0.5         0.8        0.2          0.2             0.8\n"
    "  error = 0.0            error = +0.6            error = -0.6\n"
    '  "Go straight"          "Steer left"            "Steer right"',
    top=Inches(2.2), size=Pt(15))

add_text(sl, "Why two sensors are better:", top=Inches(4.5), size=Pt(18), bold=True)
add_bullet_list(sl, [
    "One sensor:  follows the EDGE of the line, can only see one side",
    "Two sensors:  follows the CENTER of the line, can see BOTH sides",
    "No setpoint variable needed \u2014 the difference IS the error",
], top=Inches(5.0), size=Pt(17))


# ━━━ SLIDE 11: Big Picture ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "The Big Picture \u2014 Why This Matters")

add_text(sl, "Proportional control is not just a robotics trick.  It is everywhere:",
         top=Inches(0.9), size=Pt(18))

add_table(sl,
    ["System", "Setpoint", "Sensor", "Correction"],
    [
        ["Shower temperature", "Desired temp", "Your skin", "Turn the knob"],
        ["Cruise control", "Set speed", "Speedometer", "Throttle up/down"],
        ["Line-following robot", "0.5 (edge)", "Reflectance sensor", "Steer left/right"],
        ["Thermostat", "72\u00b0F", "Thermometer", "Heat on/off"],
    ],
    top=Inches(1.5), row_height=Inches(0.42), font_size=Pt(16))

add_text(sl, "The same pattern every time:", top=Inches(3.8), size=Pt(20), bold=True,
         color=ACCENT_BLUE)
add_bullet_list(sl, [
    "1.  Measure where you are",
    "2.  Calculate how far off you are  (error)",
    "3.  Respond proportionally",
], top=Inches(4.3), size=Pt(20), bold_items={0, 1, 2})

add_text(sl, "You now understand one of the most important ideas in engineering.",
         top=Inches(5.7), size=Pt(22), bold=True, color=ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)


# ━━━ SLIDE 12: Quick Check ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Quick Check \u2014 Can You Answer These?")

questions = [
    "1.  The sensor reads 0.3 and the setpoint is 0.5.  What is the error?\n"
    "     Is it positive or negative?  Which way should the robot steer?",
    "2.  Two students both have error = 0.3.  Student A uses Kp = 0.2,\n"
    "     Student B uses Kp = 1.0.  Whose robot steers harder?\n"
    "     Whose robot is more likely to oscillate?",
    "3.  Your robot is zigzagging wildly.  Should you increase or decrease Kp?",
    "4.  Your robot slowly drifts off the line on curves.\n"
    "     Should you increase or decrease Kp?",
    '5.  Why can\'t we use straight() and turn() for line following?\n'
    '     What makes arcade() different?',
]

txbox = sl.shapes.add_textbox(MARGIN_LEFT, Inches(1.0), CONTENT_W, Inches(5.5))
tf = txbox.text_frame
tf.word_wrap = True

for i, q in enumerate(questions):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = q
    run.font.name = FONT_BODY
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_GRAY

add_text(sl, "Discuss with a partner, then we will go through them together.",
         top=Inches(6.1), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/bradmiller/github/IntoToPython/module-02-line-tracking/slides/proportional-control-review.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
