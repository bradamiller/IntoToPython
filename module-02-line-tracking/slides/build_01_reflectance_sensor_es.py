"""Build the Spanish version of Lesson 1: The Reflectance Sensor slide deck.

Uses the same styling conventions as the other module-02 build scripts.
Run: /opt/homebrew/anaconda3/bin/python3 build_01_reflectance_sensor_es.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants matching existing deck styling ─────────────────────────────────
SLIDE_W = 12191695  # 13.33 in (widescreen 16:9)
SLIDE_H = 6858000   # 7.50 in

HEADER_COLOR = RGBColor(0x1B, 0x3A, 0x5C)  # dark navy
HEADER_H = 685800   # ~0.75 in
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)
CODE_BORDER = RGBColor(0x99, 0x99, 0x99)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)

MARGIN_LEFT = Inches(0.6)
CONTENT_W = Inches(12.0)

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"
FONT_TITLE = "Calibri"


# ── Helper functions ─────────────────────────────────────────────────────────

def add_header_bar(slide, height=HEADER_H):
    """Add the dark navy header bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_COLOR
    shape.line.fill.background()
    return shape


def add_title(slide, text, top=Inches(0.12), size=Pt(32)):
    """Add white title text over the header bar."""
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_W, Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TITLE
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = WHITE
    return txbox


def add_text(slide, text, left=MARGIN_LEFT, top=Inches(1.0), width=CONTENT_W,
             height=Inches(0.5), size=Pt(18), bold=False, color=DARK_GRAY,
             alignment=None, font_name=FONT_BODY):
    """Add a simple text box."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txbox


def add_bullet_list(slide, items, left=MARGIN_LEFT, top=Inches(1.5),
                    width=CONTENT_W, size=Pt(16), color=DARK_GRAY,
                    bold_items=None):
    """Add a bulleted text list. bold_items is a set of indices to bold."""
    height = Inches(0.4 * len(items))
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    bold_items = bold_items or set()

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"  \u2022  {item}"
        run.font.name = FONT_BODY
        run.font.size = size
        run.font.bold = i in bold_items
        run.font.color.rgb = color
    return txbox


def add_code_block(slide, code, left=MARGIN_LEFT, top=Inches(2.0),
                   width=Inches(11.0), size=Pt(14)):
    """Add code in a rounded-rectangle code block."""
    lines = code.strip().split("\n")
    line_h = size.pt * 1.4
    height = Emu(int((len(lines) * line_h + 20) * 12700))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CODE_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_CODE
        run.font.size = size
        run.font.color.rgb = DARK_GRAY
    return shape


def add_table(slide, headers, rows, left=MARGIN_LEFT, top=Inches(2.5),
              width=Inches(11.5), row_height=Inches(0.35), font_size=Pt(14),
              header_color=ACCENT_BLUE, alt_row_color=LIGHT_BLUE):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    col_w = int(width / num_cols)

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = tbl_shape.table

    for ci in range(num_cols):
        table.columns[ci].width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT_BODY
                run.font.size = font_size
                run.font.bold = True
                run.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_color
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = font_size
                    run.font.color.rgb = DARK_GRAY

    return tbl_shape


# ── Build the presentation ───────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank


# ━━━ SLIDE 1: Título y Objetivos de Aprendizaje ━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl, height=Emu(2743200))
add_text(sl, "Módulo 2: Seguimiento de Línea", top=Inches(0.5), size=Pt(22),
         color=WHITE, left=MARGIN_LEFT)
add_text(sl, "El Sensor de Reflectancia", top=Inches(1.2), size=Pt(36),
         bold=True, color=WHITE, left=MARGIN_LEFT)

add_text(sl, "Objetivos de aprendizaje:", top=Inches(3.0), size=Pt(20),
         bold=True, color=ACCENT_BLUE)
add_bullet_list(sl, [
    "Explicar cómo los sensores de reflectancia detectan superficies claras y oscuras",
    "Leer valores del sensor usando la clase Reflectance",
    "Registrar datos de calibración para tu robot y superficie específicos",
    "Elegir un valor de umbral para la detección de líneas",
], top=Inches(3.5), size=Pt(18))

add_text(sl, "Agenda:  Cómo funcionan los sensores (10 min)  •  Lectura de valores (15 min)  •  Ejercicio de calibración (20 min)",
         top=Inches(5.5), size=Pt(16), color=MED_GRAY)


# ━━━ SLIDE 2: Gancho — ¿Cómo "ve" un robot? ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, '¿Cómo "Ve" un Robot?')

add_text(sl, "Pregunta: \"Tú puedes ver la línea negra. ¿Cómo podría un robot notar la diferencia?\"",
         top=Inches(1.0), size=Pt(20), bold=True, color=ACCENT_BLUE)

add_text(sl, "Discusión: Los robots no tienen ojos — usan sensores que miden propiedades físicas.",
         top=Inches(2.0), size=Pt(20))

add_text(sl, "Hoy: Aprenderemos a leer los sensores de reflectancia del XRP — la forma\n"
             "en que el robot \"ve\" la línea.",
         top=Inches(3.0), size=Pt(20), bold=True)


# ━━━ SLIDE 3: Cómo Funcionan los Sensores de Reflectancia ━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Cómo Funcionan los Sensores de Reflectancia")

add_text(sl, "El sensor tiene dos partes:", top=Inches(0.9), size=Pt(20), bold=True)

add_bullet_list(sl, [
    "LED infrarrojo — emite luz hacia abajo",
    "Detector de luz — mide cuánta luz rebota de vuelta",
], top=Inches(1.4), size=Pt(18), bold_items={0, 1})

add_text(sl, "Superficie blanca: Refleja la mayoría de la luz → valor bajo (cercano a 0.0)\n"
             "Superficie negra: Absorbe la mayoría de la luz → valor alto (cercano a 1.0)",
         top=Inches(2.5), size=Pt(18))

add_text(sl, "Analogía: Una linterna sobre un espejo vs. una camiseta negra",
         top=Inches(3.5), size=Pt(18), bold=True, color=ACCENT_BLUE)

add_text(sl, "Mostrar: Foto de la parte inferior del XRP con los sensores resaltados",
         top=Inches(4.3), size=Pt(16), color=MED_GRAY)


# ━━━ SLIDE 4: Rango de Valores del Sensor ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Rango de Valores del Sensor")

add_text(sl, "Los valores van de 0.0 a 1.0:", top=Inches(0.9), size=Pt(20), bold=True)

add_table(sl,
    ["Posición", "Sensor Izquierdo", "Sensor Derecho"],
    [
        ["En superficie blanca", "~0.1 – 0.3", "~0.1 – 0.3"],
        ["En el borde de la línea", "~0.4 – 0.6", "~0.4 – 0.6"],
        ["En la línea negra", "~0.7 – 0.9", "~0.7 – 0.9"],
    ],
    top=Inches(1.5), row_height=Inches(0.42), font_size=Pt(16))

add_text(sl, "Idea clave: La transición de blanco a negro es GRADUAL, no repentina.",
         top=Inches(3.5), size=Pt(20), bold=True, color=ACCENT_BLUE)

add_text(sl, "El XRP tiene DOS sensores — izquierdo y derecho — cada uno lee de forma independiente.",
         top=Inches(4.3), size=Pt(18))


# ━━━ SLIDE 5: Lectura de Sensores en Python ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Lectura de Sensores en Python")

add_text(sl, "Importar y configurar:", top=Inches(0.9), size=Pt(18), bold=True)

add_code_block(sl,
    "from XRPLib.reflectance import Reflectance\n"
    "from XRPLib.board import Board\n"
    "\n"
    "reflectance = Reflectance.get_default_reflectance()\n"
    "board = Board.get_default_board()",
    top=Inches(1.3), size=Pt(15))

add_text(sl, "Leer valores:", top=Inches(3.3), size=Pt(18), bold=True)

add_code_block(sl,
    "board.wait_for_button()\n"
    "left = reflectance.get_left()\n"
    "right = reflectance.get_right()\n"
    'print("Izquierdo:", left, "Derecho:", right)',
    top=Inches(3.7), size=Pt(15))

add_text(sl, "Punto clave: get_left() y get_right() devuelven un número entre 0.0 y 1.0.",
         top=Inches(5.5), size=Pt(18), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 6: Lectura Continua ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Lectura Continua")

add_text(sl, "Leer muchas veces en un bucle:", top=Inches(0.9), size=Pt(18), bold=True)

add_code_block(sl,
    "import time\n"
    "\n"
    "board.wait_for_button()\n"
    "\n"
    "for i in range(50):\n"
    "    left = reflectance.get_left()\n"
    "    right = reflectance.get_right()\n"
    '    print("Izquierdo:", left, "  Derecho:", right)\n'
    "    time.sleep(0.1)",
    top=Inches(1.4), size=Pt(15))

add_text(sl, "¿Por qué? Para ver cómo cambian los valores al mover el robot sobre la línea.",
         top=Inches(4.2), size=Pt(18), bold=True)

add_text(sl, "time.sleep(0.1) — pausa 0.1 segundos para que las lecturas no sean demasiado rápidas.",
         top=Inches(4.8), size=Pt(16), color=MED_GRAY, font_name=FONT_CODE)


# ━━━ SLIDE 7: ¿Qué es un Umbral? ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "¿Qué Es un Umbral?")

add_text(sl, 'Umbral: Un valor de corte para decidir "sobre la línea" vs. "fuera de la línea"',
         top=Inches(0.9), size=Pt(20), bold=True)

add_text(sl, "Ejemplo:", top=Inches(1.6), size=Pt(18), bold=True, color=ACCENT_BLUE)

add_bullet_list(sl, [
    "Lecturas fuera de la línea: 0.1, 0.2, 0.15, 0.25",
    "Lecturas sobre la línea: 0.75, 0.8, 0.85, 0.9",
    "Umbral = 0.5 (punto medio entre ambos)",
], top=Inches(2.0), size=Pt(18))

add_text(sl, "Regla: Si sensor > umbral → sobre la línea.  Si sensor < umbral → fuera de la línea.",
         top=Inches(3.5), size=Pt(20), bold=True, color=ACCENT_BLUE)

add_text(sl, "¡Importante! El umbral depende de TU superficie y cinta. ¡Hay que calibrar!",
         top=Inches(4.3), size=Pt(20), bold=True)


# ━━━ SLIDE 8: Ejercicio de Calibración ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Ejercicio de Calibración")

add_text(sl, "Tu tarea:", top=Inches(0.9), size=Pt(20), bold=True)

add_bullet_list(sl, [
    "Coloca el robot en superficie BLANCA → registra los valores del sensor",
    "Coloca el robot en la línea NEGRA → registra los valores del sensor",
    "Coloca el robot en el BORDE de la línea → registra los valores del sensor",
    "Calcula tu umbral (punto medio entre los promedios fuera y sobre la línea)",
], top=Inches(1.5), size=Pt(18))

add_text(sl, "¿Por qué calibrar? Diferentes superficies, colores de cinta y condiciones\n"
             "de iluminación dan valores diferentes.",
         top=Inches(3.5), size=Pt(18))

add_text(sl, "¡Guarda tu umbral! Lo usarás durante el resto del Módulo 2.",
         top=Inches(4.5), size=Pt(20), bold=True, color=ACCENT_BLUE)


# ━━━ SLIDE 9: ¡Tu Turno! ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "¡Tu Turno!")

add_text(sl, "Actividad:", top=Inches(0.9), size=Pt(20), bold=True)

add_bullet_list(sl, [
    "Crea sensor_test.py",
    "Escribe el programa de lectura continua",
    "Recoge datos de calibración para 3 posiciones",
    "Regístralos en tu hoja de trabajo",
    "Determina tu valor de umbral",
], top=Inches(1.4), size=Pt(18))

add_text(sl, "Puntos de verificación:", top=Inches(3.8), size=Pt(20), bold=True,
         color=ACCENT_BLUE)

add_bullet_list(sl, [
    "¿Puedes leer los valores del sensor?",
    "¿Dan números diferentes el blanco y el negro?",
    "¿Encontraste un buen umbral?",
], top=Inches(4.3), size=Pt(18))


# ━━━ SLIDE 10: Conexión con la Siguiente Lección ━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(blank_layout)
add_header_bar(sl)
add_title(sl, "Conexión con la Siguiente Lección")

add_text(sl, "Lo que hiciste hoy:", top=Inches(0.9), size=Pt(20), bold=True)

add_bullet_list(sl, [
    "Aprendiste cómo funcionan los sensores de reflectancia",
    "Leíste valores del sensor en Python",
    "Calibraste tus sensores y elegiste un umbral",
], top=Inches(1.4), size=Pt(18))

add_text(sl, "Próxima lección (Lección 2):", top=Inches(3.0), size=Pt(20), bold=True,
         color=ACCENT_BLUE)

add_bullet_list(sl, [
    "Usar un bucle while para mantener el robot en movimiento",
    "Detenerse cuando el sensor detecte la línea",
    "¡Primera vez que el robot usa sensores para controlar su comportamiento!",
], top=Inches(3.5), size=Pt(18))

add_text(sl, "¡Guarda tus datos de calibración — los vas a necesitar!",
         top=Inches(5.0), size=Pt(22), bold=True, color=ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/bradmiller/GitHub/IntoToPython/module-02-line-tracking/slides/01-the-reflectance-sensor-es.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
