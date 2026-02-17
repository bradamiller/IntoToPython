#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 7: Detecting Intersections."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              row_height=Inches(0.45), font_size=Pt(14)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        cw = int(width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = cw
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.0))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 7: Detecting Intersections"
    p.font.size = Pt(42)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for obj in [
        "Detect when both sensors are on the line simultaneously",
        "Use the and logical operator for combined conditions",
        "Integrate intersection detection with line following",
        "Program the robot to reverse direction at a cross",
    ]:
        add_bullet_text(tf, obj, level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5.5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for item, duration in [
        ("What intersections look like to sensors", "10 min"),
        ("Cross detection code", "15 min"),
        ("Follow + detect + reverse", "20 min"),
    ]:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: Adding the Cross
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Adding the Cross")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "New physical setup:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "A taped cross \u2014 a perpendicular line crossing the circle.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Question:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, '"When the robot reaches this cross while following the line, what will the sensors see?"',
                    level=0, font_size=Pt(20), color=ACCENT_BLUE, bold=True)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Answer:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "BOTH sensors will be on the line at the same time!",
                    level=0, font_size=Pt(20), bold=True, color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 3: How Sensors See an Intersection
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "How Sensors See an Intersection")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "During normal line following:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "One sensor on line, one off \u2192 different values",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "left \u2248 0.8,  right \u2248 0.2  (or vice versa)",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(6.0), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "At an intersection (cross):", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "BOTH sensors on the line \u2192 both read high",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "left > 0.5  AND  right > 0.5",
                    level=0, font_size=Pt(18), bold=True, color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "This is our detection signal!"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 4: The and Operator
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The and Operator")

    add_code_block(slide,
        "if left > 0.5 and right > 0.5:\n"
        '    print("CROSS DETECTED!")',
        Inches(0.6), Inches(1.5), Inches(8.0), Inches(1.3)
    )

    add_label(slide, "and truth table:", Inches(0.6), Inches(3.1), Inches(12), Inches(0.5))

    headers = ["left > 0.5", "right > 0.5", "Result"]
    rows = [
        ["True",  "True",  "True \u2190 Cross!"],
        ["True",  "False", "False"],
        ["False", "True",  "False"],
        ["False", "False", "False"],
    ]
    col_widths = [Inches(3.0), Inches(3.0), Inches(4.5)]
    add_table(slide, headers, rows, Inches(0.6), Inches(3.7), Inches(10.0),
              col_widths=col_widths, row_height=Inches(0.48))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Both must be True for and to be True."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 5: Line Following + Cross Detection
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Line Following + Cross Detection")

    add_code_block(slide,
        "while True:\n"
        "    left  = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "\n"
        "    if left > threshold and right > threshold:\n"
        "        # Cross detected!\n"
        "        drivetrain.stop()\n"
        '        print("Cross!")\n'
        "        drivetrain.turn(180)\n"
        "        time.sleep(0.3)\n"
        "    else:\n"
        "        # Normal line following\n"
        "        error        = left - right\n"
        "        left_effort  = base_effort - error * Kp\n"
        "        right_effort = base_effort + error * Kp\n"
        "        drivetrain.set_effort(left_effort, right_effort)",
        Inches(0.6), Inches(1.4), Inches(8.5), Inches(5.8),
        font_size=Pt(13)
    )

    txBox = slide.shapes.add_textbox(Inches(9.2), Inches(1.4), Inches(3.9), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Key design:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Check for cross FIRST, then do normal following.",
                    level=0, font_size=Pt(17), bold=True, color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Cross detected:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Stop, print, turn 180\u00b0, wait", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Not a cross:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Two-sensor proportional following", level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 6: Why time.sleep() After Turning?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Why time.sleep() After Turning?")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Without sleep:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot turns 180\u00b0 \u2192 immediately re-reads sensors \u2192 still at the cross \u2192 turns again!",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "With sleep(0.3):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot turns 180\u00b0 \u2192 pauses briefly \u2192 moves past the cross \u2192 resumes following",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "The sleep gives the robot time to move away from the intersection.",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 7: Counting Crossings
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Counting Crossings")

    add_code_block(slide,
        "cross_count = 0\n"
        "\n"
        "while cross_count < 4:\n"
        "    left  = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "\n"
        "    if left > threshold and right > threshold:\n"
        "        cross_count = cross_count + 1\n"
        "        drivetrain.stop()\n"
        '        print("Cross #", cross_count)\n'
        "        drivetrain.turn(180)\n"
        "        time.sleep(0.3)\n"
        "    else:\n"
        "        # line following code...\n"
        "\n"
        'print("Done! Crossed", cross_count, "times.")',
        Inches(0.6), Inches(1.4), Inches(9.5), Inches(5.5)
    )

    txBox = slide.shapes.add_textbox(Inches(10.2), Inches(1.4), Inches(2.9), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "This previews the final project!",
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 8: Your Turn!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Your Turn!")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Exercise 1:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Follow the circle and reverse at the cross", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Use two-sensor line following from Lesson 6", level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Add cross detection with and", level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Turn 180\u00b0 when cross is detected", level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Exercise 2 (Challenge):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Count crossings and stop after 4", level=0, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.0), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Test:", level=0, font_size=Pt(22), bold=True, color=ACCENT_BLUE)
    add_bullet_text(tf, "Does the robot reliably detect the cross?",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 9: Connection to Next Lesson
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Next Lesson")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "What you\u2019ve built so far (all as loose code):",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Sensor reading, error calculation, cross detection",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Proportional control, differential steering",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "All mixed together in one big program",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next (Lesson 8): Package sensor logic into a LineSensor class",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "get_error(),  is_at_cross(),  is_off_line()",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "First time writing a Python class!",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Why? Organized code is easier to understand, test, and reuse.",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "07-detecting-intersections.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
