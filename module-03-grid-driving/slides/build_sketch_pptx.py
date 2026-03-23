#!/usr/bin/env python3
"""Build a sketch-style PPTX for Module 3 Lesson 1 using rough.js illustrations."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Sketch Color Palette ──
TEAL = RGBColor(0x5B, 0x8F, 0xA8)
ORANGE = RGBColor(0xE8, 0xA8, 0x7C)
GREEN = RGBColor(0x85, 0xC8, 0x8A)
DARK = RGBColor(0x2D, 0x34, 0x36)
MUTED = RGBColor(0x63, 0x6E, 0x72)
BORDER_GRAY = RGBColor(0xB8, 0xC5, 0xD0)
PAPER_BG = RGBColor(0xFD, 0xFC, 0xFA)
CODE_BG = RGBColor(0xF0, 0xED, 0xE8)
LIGHT_TEAL_BG = RGBColor(0xE4, 0xEF, 0xF4)
LIGHT_ORANGE_BG = RGBColor(0xFD, 0xF0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Fonts ──
TITLE_FONT = "Caveat"
BODY_FONT = "Patrick Hand"
CODE_FONT = "Fira Code"

# ── Dimensions ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ILLUST_DIR = os.path.join(os.path.dirname(__file__), '_illustrations')


def img(name):
    return os.path.join(ILLUST_DIR, name)


def set_slide_bg(slide, color=PAPER_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_sketch_title_bar(slide, text):
    """Light title bar with teal bottom border."""
    # Light background strip
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_TEAL_BG
    bar.line.fill.background()

    # Bottom border line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.1),
        SLIDE_W, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_text(slide, text, left, top, width, height,
             font_name=BODY_FONT, font_size=Pt(20), color=DARK,
             bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_header(slide, text, left, top, width=Inches(10)):
    return add_text(slide, text, left, top, width, Inches(0.5),
                    font_name=TITLE_FONT, font_size=Pt(28),
                    color=TEAL, bold=True)


def add_bullet(tf, text, level=0, font_size=Pt(20), color=DARK, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(4)
    return p


def add_code_block(slide, code, left, top, width, height):
    """Code block with sketch-style background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.line.color.rgb = BORDER_GRAY
    bg.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.15),
        width - Inches(0.5), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK
        p.space_after = Pt(2)


def add_sketch_box(slide, left, top, width, height, border_color=TEAL):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PAPER_BG
    box.line.color.rgb = border_color
    box.line.width = Pt(2)
    return box


def add_insight_box(slide, text, left, top, width):
    box = add_sketch_box(slide, left, top, width, Inches(0.7), ORANGE)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_ORANGE_BG
    add_text(slide, text, left + Inches(0.2), top + Inches(0.1),
             width - Inches(0.4), Inches(0.5),
             font_name=TITLE_FONT, font_size=Pt(24), color=DARK,
             bold=True, alignment=PP_ALIGN.CENTER)


def add_table(slide, headers, rows, left, top, width):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    row_h = Inches(0.45)
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_h * num_rows
    )
    table = tbl_shape.table
    col_w = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_w

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_TEAL_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = TEAL
            p.font.name = TITLE_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# BUILD SLIDES
# ═══════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: Title ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)

    # Robot illustration
    slide.shapes.add_picture(img('illust-robot.png'),
                             Inches(5.8), Inches(0.3), Inches(1.6))

    # Module label
    add_text(slide, 'MODULE 3  ·  LESSON 1',
             Inches(2), Inches(1.8), Inches(9), Inches(0.5),
             font_name=BODY_FONT, font_size=Pt(22), color=TEAL,
             alignment=PP_ALIGN.CENTER)

    # Title
    add_text(slide, 'Introduction to the Grid',
             Inches(1.5), Inches(2.3), Inches(10), Inches(1.0),
             font_name=TITLE_FONT, font_size=Pt(54), color=DARK,
             bold=True, alignment=PP_ALIGN.CENTER)

    # Underline accent
    ul = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), Inches(3.3), Inches(4), Pt(4)
    )
    ul.fill.solid()
    ul.fill.fore_color.rgb = ORANGE
    ul.line.fill.background()

    # Learning Objectives
    add_text(slide, 'Learning Objectives',
             Inches(1.2), Inches(3.8), Inches(5), Inches(0.5),
             font_name=TITLE_FONT, font_size=Pt(26), color=TEAL, bold=True)
    tf = add_text(slide, '', Inches(1.2), Inches(4.3), Inches(5.5), Inches(3.0))
    tf.paragraphs[0].text = ''
    for obj in [
        'Describe the physical grid layout and how intersections are formed',
        'Explain how cross detection maps to grid intersections',
        'Use LineTrack to drive to the first intersection',
        'Reuse existing code on a new surface',
    ]:
        add_bullet(tf, obj, font_size=Pt(18), color=DARK)

    # Agenda
    add_text(slide, 'Agenda',
             Inches(7.5), Inches(3.8), Inches(5), Inches(0.5),
             font_name=TITLE_FONT, font_size=Pt(26), color=TEAL, bold=True)
    tf = add_text(slide, '', Inches(7.5), Inches(4.3), Inches(5), Inches(3.0))
    tf.paragraphs[0].text = ''
    for item in [
        'From circle to grid (10 min)',
        'Review: LineTrack methods (5 min)',
        'Guided: Drive to first intersection (15 min)',
        'Practice (15 min)',
    ]:
        add_bullet(tf, item, font_size=Pt(18), color=DARK)

    # ── SLIDE 2: From Circle to Grid ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'From Circle to Grid')

    add_header(slide, 'Remember Module 2?', Inches(0.8), Inches(1.3))
    tf = add_text(slide, 'You built a LineTrack class that follows lines and detects crosses',
                  Inches(0.8), Inches(1.8), Inches(11), Inches(0.4))

    # Illustration
    slide.shapes.add_picture(img('illust-circle-to-grid.png'),
                             Inches(2.5), Inches(2.4), Inches(8))

    # Question box
    add_sketch_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6), ORANGE)
    add_text(slide, '"Can your Module 2 code work on this grid without changes?"',
             Inches(1.8), Inches(5.25), Inches(9.5), Inches(0.5),
             font_size=Pt(20), bold=True)

    # Answer box
    add_sketch_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.6), GREEN)
    add_text(slide, "Yes! That's the power of code reuse.",
             Inches(1.8), Inches(5.95), Inches(9.5), Inches(0.5),
             font_size=Pt(20), bold=True, color=GREEN)

    # ── SLIDE 3: The Grid Layout ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'The Grid Layout')

    add_header(slide, 'A grid is a network of taped lines', Inches(0.8), Inches(1.3))

    # Grid illustration
    slide.shapes.add_picture(img('illust-grid-layout.png'),
                             Inches(3.5), Inches(1.8), Inches(6.5))

    # Box with key points
    add_sketch_box(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(1.5), TEAL)
    add_header(slide, 'Each orange dot is an intersection', Inches(1.3), Inches(5.7), Inches(10))
    tf = add_text(slide, '', Inches(1.3), Inches(6.1), Inches(10), Inches(0.9))
    tf.paragraphs[0].text = ''
    add_bullet(tf, 'Lines connect them horizontally and vertically')
    add_bullet(tf, 'The robot follows lines between intersections')
    add_bullet(tf, 'Cross detection tells the robot "you\'ve arrived!"')

    # ── SLIDE 4: Why the Grid Matters ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Why the Grid Matters')

    add_header(slide, 'The grid is a stepping stone', Inches(0.8), Inches(1.3))

    # Progression illustration
    slide.shapes.add_picture(img('illust-progression.png'),
                             Inches(2.0), Inches(2.0), Inches(9))

    # Puzzle pieces illustration
    slide.shapes.add_picture(img('illust-puzzle-pieces.png'),
                             Inches(3.0), Inches(4.8), Inches(7))

    # Insight
    add_insight_box(slide, 'Grid navigation = line following + cross detection + turns',
                    Inches(1.5), Inches(3.9), Inches(10))

    # Subtext
    add_text(slide, 'You already have all three from Module 2!',
             Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
             font_name=TITLE_FONT, font_size=Pt(22), color=TEAL,
             alignment=PP_ALIGN.CENTER)

    # ── SLIDE 5: LineTrack Methods ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Review — LineTrack Methods')

    add_header(slide, 'Your LineTrack class has three key methods:', Inches(0.8), Inches(1.3))

    add_code_block(slide, '''tracker = LineTrack()

# Follow the line until an intersection
tracker.track_until_cross()

# Turn right onto perpendicular line
tracker.turn_right()

# Turn left onto perpendicular line
tracker.turn_left()''',
                   Inches(0.8), Inches(1.9), Inches(6), Inches(3.2))

    # Method cards illustration
    slide.shapes.add_picture(img('illust-method-cards.png'),
                             Inches(7.2), Inches(2.0), Inches(5.5))

    # Insight box
    add_insight_box(slide, 'These three methods are ALL you need for grid navigation!',
                    Inches(1.5), Inches(5.5), Inches(10))

    # ── SLIDE 6: Driving to First Intersection ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Driving to the First Intersection')

    add_header(slide, 'Setup: Place robot on a grid line, behind an intersection',
               Inches(0.8), Inches(1.3), Inches(12))

    # Code
    add_code_block(slide, '''board = Board.get_default_board()
tracker = LineTrack()

board.wait_for_button()
print("Driving to first intersection...")
tracker.track_until_cross()
print("Intersection reached!")''',
                   Inches(0.8), Inches(2.0), Inches(6), Inches(2.8))

    # Illustration
    add_header(slide, 'What happens:', Inches(7.2), Inches(2.0))
    slide.shapes.add_picture(img('illust-drive-to-intersection.png'),
                             Inches(7.0), Inches(2.6), Inches(5.5))

    # Steps
    tf = add_text(slide, '', Inches(0.8), Inches(5.2), Inches(11), Inches(1.5))
    tf.paragraphs[0].text = ''
    add_bullet(tf, '1. Robot follows the line using proportional control', font_size=Pt(19))
    add_bullet(tf, '2. When both sensors detect the cross → stops', font_size=Pt(19))
    add_bullet(tf, '3. Robot is now AT the intersection', font_size=Pt(19))

    # ── SLIDE 7: Drive and Turn ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Drive and Turn')

    add_header(slide, 'After reaching an intersection, you can turn:', Inches(0.8), Inches(1.3))

    # Code
    add_code_block(slide, '''board.wait_for_button()

tracker.track_until_cross()
print("At intersection!")

tracker.turn_right()
print("Turned right!")

tracker.track_until_cross()
print("At next intersection!")''',
                   Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.5))

    # Illustration
    slide.shapes.add_picture(img('illust-drive-and-turn.png'),
                             Inches(7.0), Inches(1.8), Inches(5))

    # ── SLIDE 8: Code Reuse ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Code Reuse in Action')

    # Arrow illustration
    slide.shapes.add_picture(img('illust-code-reuse-arrow.png'),
                             Inches(2.5), Inches(1.3), Inches(8.5))

    # Table
    add_table(slide,
              ['Module 2', 'Module 3'],
              [
                  ['Follow circle', 'Follow grid lines'],
                  ['Detect cross on circle', 'Detect grid intersection'],
                  ['Turn at cross', 'Turn at intersection'],
                  ['track_until_cross()', 'Same method, same code!'],
                  ['turn_right()', 'Same method, same code!'],
              ],
              Inches(1.5), Inches(2.8), Inches(10))

    # Insight
    add_insight_box(slide, 'Good code works in new situations without modification.',
                    Inches(2), Inches(5.8), Inches(9))

    # ── SLIDE 9: Your Turn! ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, 'Your Turn!')

    add_header(slide, 'Activity', Inches(0.8), Inches(1.3))

    # Steps in a sketch box
    add_sketch_box(slide, Inches(0.6), Inches(1.8), Inches(7), Inches(3.5), TEAL)
    steps = [
        '1.  Copy your LineSensor and LineTrack classes from Module 2',
        '2.  Write a program that drives to the first intersection and stops',
        '3.  Verify the robot stops at the right place',
        '4.  Add a turn_right() after reaching the intersection',
        '5.  Add another track_until_cross() to drive to the next intersection',
    ]
    tf = add_text(slide, '', Inches(0.9), Inches(2.0), Inches(6.5), Inches(3.0))
    tf.paragraphs[0].text = ''
    for step in steps:
        add_bullet(tf, step, font_size=Pt(19))

    # Checkpoints
    add_header(slide, 'Checkpoints', Inches(8.0), Inches(1.3))
    add_sketch_box(slide, Inches(7.8), Inches(1.8), Inches(5), Inches(3.5), GREEN)
    checks = [
        'Does the robot stop at the first intersection?',
        'Does the turn line up with the perpendicular line?',
        'Does the robot reach the second intersection after turning?',
    ]
    tf = add_text(slide, '', Inches(8.1), Inches(2.0), Inches(4.5), Inches(3.0))
    tf.paragraphs[0].text = ''
    for check in checks:
        p = add_bullet(tf, check, font_size=Pt(18))

    # ── SLIDE 10: What's Next ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_sketch_title_bar(slide, "What's Next")

    # Today column
    add_header(slide, 'What you did today', Inches(0.8), Inches(1.3))
    add_sketch_box(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(2.2), TEAL)
    tf = add_text(slide, '', Inches(0.9), Inches(1.9), Inches(5), Inches(2.0))
    tf.paragraphs[0].text = ''
    add_bullet(tf, "Saw how the grid relates to Module 2's cross detection")
    add_bullet(tf, 'Drove to the first intersection using LineTrack')
    add_bullet(tf, 'Added a turn at the intersection')

    # Next lesson column
    add_header(slide, 'Next lesson (Lesson 2)', Inches(7.0), Inches(1.3))
    add_sketch_box(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(2.2), ORANGE)
    tf = add_text(slide, '', Inches(7.1), Inches(1.9), Inches(5.5), Inches(2.0))
    tf.paragraphs[0].text = ''
    add_bullet(tf, 'Drive PAST an intersection to reach the next one')
    add_bullet(tf, 'Learn the "clear the intersection" technique')
    add_bullet(tf, 'Use for loops to drive multiple intersections')

    # Challenge illustration
    slide.shapes.add_picture(img('illust-three-intersections.png'),
                             Inches(2.0), Inches(4.5), Inches(9))

    # Challenge question
    add_sketch_box(slide, Inches(2.5), Inches(5.8), Inches(8), Inches(0.8), ORANGE)
    add_text(slide, 'Challenge: What if you need to go 3 intersections forward? That\'s next!',
             Inches(2.8), Inches(5.9), Inches(7.5), Inches(0.6),
             font_name=TITLE_FONT, font_size=Pt(24), color=DARK,
             bold=True, alignment=PP_ALIGN.CENTER)

    return prs


if __name__ == '__main__':
    prs = build_presentation()
    output = os.path.join(os.path.dirname(__file__),
                          '01-introduction-to-the-grid-sketch-test.pptx')
    prs.save(output)
    print(f"Saved: {output}")
