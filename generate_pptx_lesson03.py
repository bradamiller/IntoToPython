#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 3: Bounce Driving."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)

# ── Fonts ──
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

# ── Sizes ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5),
        Inches(12), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(3.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 3: Bounce Driving"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    objectives = [
        "Write if/else statements in Python",
        "Use while True for continuous robot behavior",
        "Combine sensor checks with decision-making",
        "Program the robot to bounce between circle edges",
    ]
    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    agenda = [
        ("if/else explained", "10 min"),
        ("Bounce driving walkthrough", "15 min"),
        ("Practice and experiment", "20 min"),
    ]
    for item, duration in agenda:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: From Stopping to Bouncing
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Stopping to Bouncing")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Lesson 2:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drives to edge \u2192 STOPS.  Done.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Today:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drives to edge \u2192 TURNS AROUND \u2192 drives again \u2192 repeat forever!",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Like a ball bouncing off walls.", level=0, font_size=Pt(22), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 3: if/else Syntax
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "if/else Syntax")

    # Syntax code block (left)
    add_label(slide, "Syntax:", Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.5))
    add_code_block(slide,
        "if condition:\n"
        "    # runs when condition is True\n"
        "else:\n"
        "    # runs when condition is False",
        Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.8)
    )

    # Example code block (right)
    add_label(slide, "Example:", Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5))
    add_code_block(slide,
        "temperature = 85\n"
        "if temperature > 80:\n"
        '    print("It\'s hot!")\n'
        "else:\n"
        '    print("It\'s comfortable.")',
        Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.0)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Only ONE block runs \u2014 never both."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 4: if/else for Robot Decisions
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "if/else for Robot Decisions")

    add_code_block(slide,
        "left = reflectance.get_left()\n"
        "\n"
        "if left > threshold:\n"
        "    # ON the line \u2014 turn around!\n"
        "    drivetrain.stop()\n"
        "    drivetrain.turn(180)\n"
        "else:\n"
        "    # OFF the line \u2014 keep driving\n"
        "    drivetrain.set_effort(0.3, 0.3)",
        Inches(0.6), Inches(1.5), Inches(7.0), Inches(4.0)
    )

    txBox = slide.shapes.add_textbox(Inches(8.0), Inches(1.5), Inches(4.8), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "When sensor > threshold:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "On the line \u2192 stop and turn 180\u00b0", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "When sensor \u2264 threshold:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Off the line \u2192 keep driving", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "The robot checks the sensor and decides what to do.",
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 5: while True — The Infinite Loop
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "while True \u2014 The Infinite Loop")

    add_code_block(slide,
        "while True:\n"
        "    # this code runs FOREVER",
        Inches(0.6), Inches(1.5), Inches(7.0), Inches(1.5)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(12), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Why?", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "True is always True \u2014 the condition never becomes False.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "For robots:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "The robot should keep running until you turn it off.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Common in robotics \u2014 the robot continuously senses and acts.",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 6: Complete Bounce Program
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Complete Bounce Program")

    add_code_block(slide,
        "threshold = 0.5\n"
        "\n"
        "board.wait_for_button()\n"
        "\n"
        "while True:\n"
        "    left = reflectance.get_left()\n"
        "\n"
        "    if left > threshold:\n"
        "        drivetrain.stop()\n"
        '        print("Bounce!")\n'
        "        drivetrain.turn(180)\n"
        "    else:\n"
        "        drivetrain.set_effort(0.3, 0.3)",
        Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.2)
    )

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Flow:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Drive forward", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Check sensor", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "On line? \u2192 stop, print, turn 180\u00b0", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Not on line? \u2192 keep driving", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Repeat forever", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 7: Logical Operators
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Logical Operators")

    # or
    add_label(slide, "or \u2014 at least one must be True:",
              Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    add_code_block(slide,
        "if left > 0.5 or right > 0.5:\n"
        '    print("At least one sensor on the line!")',
        Inches(0.6), Inches(2.0), Inches(12), Inches(1.3)
    )

    # and
    add_label(slide, "and \u2014 both must be True:",
              Inches(0.6), Inches(3.5), Inches(12), Inches(0.5))
    add_code_block(slide,
        "if left > 0.5 and right > 0.5:\n"
        '    print("BOTH sensors on the line!")',
        Inches(0.6), Inches(4.1), Inches(12), Inches(1.3)
    )

    # not
    add_label(slide, "not \u2014 reverses True/False:",
              Inches(0.6), Inches(5.6), Inches(12), Inches(0.5))
    add_code_block(slide,
        "if not (left > 0.5):\n"
        '    print("Left is NOT on the line")',
        Inches(0.6), Inches(6.2), Inches(12), Inches(1.0)
    )

    # ════════════════════════════════════════
    # SLIDE 8: Using Both Sensors
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Using Both Sensors")

    add_code_block(slide,
        "while True:\n"
        "    left = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "\n"
        "    if left > threshold or right > threshold:\n"
        "        drivetrain.stop()\n"
        "        drivetrain.turn(180)\n"
        "    else:\n"
        "        drivetrain.set_effort(0.3, 0.3)",
        Inches(0.6), Inches(1.5), Inches(7.0), Inches(4.2)
    )

    txBox = slide.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(5.0), Inches(4.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Why both sensors?", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "If the robot approaches the line at an angle, only one sensor may detect it first.",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Using or catches both cases.",
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 9: Your Turn!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Your Turn!")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Exercise 1:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Build the bounce driving program", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Robot bounces back and forth inside the circle", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Exercise 2 (Challenge):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Use both sensors with or", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Add a bounce counter", level=0, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Observe:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "What pattern does the robot follow?", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Does it always bounce on the same line?", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 10: Connection to Next Lesson
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Next Lesson")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "The problem:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "The robot bounces back and forth on the SAME line.  Boring!",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next lesson (Lesson 4):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Add randomness!", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "import random", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "random.randint() \u2014 generate random numbers", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Random turn angles \u2192 robot explores the whole circle", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Preview:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)

    add_code_block(slide,
        "angle = random.randint(90, 270)\n"
        "drivetrain.turn(angle)",
        Inches(0.6), Inches(5.6), Inches(7.0), Inches(1.5)
    )

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "03-bounce-driving.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
