#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 4: Random Turns."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

# ── Fonts ──
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

# ── Sizes ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = table_shape.table

    col_widths = [int(width * 0.35), int(width * 0.65)]
    for i in range(num_cols):
        table.columns[i].width = col_widths[i] if i < len(col_widths) else int(width / num_cols)

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5),
        Inches(12), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(3.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 4: Random Turns"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    objectives = [
        "Use import to bring in Python\u2019s random module",
        "Generate random numbers with random.randint()",
        "Explain why predictable behavior is a problem",
        "Modify bounce driving to use random turn angles",
    ]
    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    agenda = [
        ("The predictability problem", "5 min"),
        ("import and random module", "10 min"),
        ("Random bounce driving", "15 min"),
        ("Experimentation", "15 min"),
    ]
    for item, duration in agenda:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: The Predictability Problem
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Predictability Problem")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Observe the Lesson 3 robot:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Drives forward \u2192 hits line \u2192 turns 180\u00b0 \u2192 drives back \u2192 hits line \u2192 turns 180\u00b0\u2026",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "Same path every time!  Never explores the circle.",
                    level=0, font_size=Pt(20), bold=True, color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Why?", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "180\u00b0 always sends it straight back.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Fix:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Turn a DIFFERENT amount each time \u2192 explore the whole area.",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 3: import — Using Python Modules
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "import \u2014 Using Python Modules")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Python has thousands of built-in tools organized into modules.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "To use one:  import module_name",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)

    add_code_block(slide,
        "import random      # Makes random functions available\n"
        "import time        # Makes time functions available",
        Inches(0.6), Inches(3.0), Inches(8.0), Inches(1.3)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "You\u2019ve already done this:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "from XRPLib.reflectance import Reflectance",
        Inches(0.6), Inches(5.1), Inches(8.0), Inches(0.9)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Rule: Import at the TOP of your file."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 4: random.randint()
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "random.randint()")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "random.randint(a, b) \u2014 returns a random integer between a and b (inclusive)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "import random\n"
        "\n"
        "random.randint(1, 6)     # Like rolling a die: 1, 2, 3, 4, 5, or 6\n"
        "random.randint(90, 270)  # Random angle for turning\n"
        "random.randint(1, 100)   # Random percentage",
        Inches(0.6), Inches(2.1), Inches(8.5), Inches(2.2)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Each call gives a DIFFERENT number:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    add_code_block(slide,
        "print(random.randint(1, 10))  # Maybe 7\n"
        "print(random.randint(1, 10))  # Maybe 3\n"
        "print(random.randint(1, 10))  # Maybe 9",
        Inches(0.6), Inches(5.1), Inches(8.5), Inches(1.7)
    )

    # ════════════════════════════════════════
    # SLIDE 5: Adding Randomness to Bounce Driving
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Adding Randomness to Bounce Driving")

    add_label(slide, "Before (predictable):", Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "drivetrain.turn(180)  # Always 180\u00b0",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.0)
    )

    add_label(slide, "After (random):", Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5))
    add_code_block(slide,
        "angle = random.randint(90, 270)\n"
        'print("Turning", angle, "degrees")\n'
        "drivetrain.turn(angle)",
        Inches(7.0), Inches(2.0), Inches(6.0), Inches(1.7)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Result: Robot turns a different amount each bounce \u2192 explores the circle!"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 6: Complete Random Bounce Program
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Complete Random Bounce Program")

    add_code_block(slide,
        "from XRPLib.reflectance import Reflectance\n"
        "from XRPLib.differential_drive import DifferentialDrive\n"
        "from XRPLib.board import Board\n"
        "import random\n"
        "\n"
        "reflectance = Reflectance.get_default_reflectance()\n"
        "drivetrain = DifferentialDrive.get_default_differential_drive()\n"
        "board = Board.get_default_board()\n"
        "\n"
        "threshold = 0.5\n"
        "board.wait_for_button()\n"
        "\n"
        "while True:\n"
        "    left = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "\n"
        "    if left > threshold or right > threshold:\n"
        "        drivetrain.stop()\n"
        "        angle = random.randint(90, 270)\n"
        '        print("Turning", angle, "degrees")\n'
        "        drivetrain.turn(angle)\n"
        "    else:\n"
        "        drivetrain.set_effort(0.3, 0.3)",
        Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.8),
        font_size=Pt(13)
    )

    # ════════════════════════════════════════
    # SLIDE 7: Experimenting with Ranges
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Experimenting with Ranges")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Try different ranges and observe:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    headers = ["Range", "Behavior"]
    rows = [
        ["randint(170, 190)", "Small variation around 180\u00b0 \u2014 mostly back and forth"],
        ["randint(90, 270)",  "Wide range \u2014 good exploration"],
        ["randint(45, 315)",  "Very wide \u2014 may turn too little sometimes"],
        ["randint(90, 90)",   "Always 90\u00b0 \u2014 not random at all!"],
    ]
    add_table(slide, headers, rows, Inches(0.8), Inches(2.1), Inches(11.5), row_height=Inches(0.5))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Question: Which range keeps the robot inside the circle best?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 8: Your Turn!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Your Turn!")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Exercise 1:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Add random turns to your bounce program", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Change drivetrain.turn(180) to use random.randint()",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Exercise 2 (Challenge):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Experiment with different ranges", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Record which range works best", level=1, font_size=Pt(18), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Bonus:", level=0, font_size=Pt(22), bold=True, color=ACCENT_BLUE)
    add_bullet_text(tf, "Randomize the speed too!", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 9: Connection to Next Lesson
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Next Lesson")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Module 2 so far:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Lesson 1: Read sensors \u2713", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Lesson 2: while loops \u2014 drive to edge \u2713", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Lesson 3: if/else \u2014 bounce driving \u2713", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Lesson 4: import random \u2014 random turns \u2713", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next (Lesson 5):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Instead of BOUNCING off the line, FOLLOW it!",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "Proportional control \u2014 smooth, continuous line following",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "The robot adjusts steering based on how far it is from the line",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Big shift: From bouncing to tracking!",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "04-random-turns.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
