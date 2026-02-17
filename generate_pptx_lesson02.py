#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 2: Drive to the Edge and Stop."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
GREEN = RGBColor(0x17, 0x7B, 0x41)
RED = RGBColor(0xC0, 0x00, 0x00)

# ── Fonts ──
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

# ── Sizes ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    """Add a colored title bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    """Add a bullet point to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    """Add a code block with gray background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height, font_size=Pt(18),
              bold=True, color=DARK_BLUE):
    """Add a simple labeled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45)):
    """Add a formatted table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = table_shape.table

    col_width = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    """Create a text box for main content below the title bar."""
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5),
        Inches(12), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # Large blue title area
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(3.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Module subtitle
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 2: Drive to the Edge and Stop"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    # Learning objectives
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    objectives = [
        "Write while loops in Python",
        "Use comparison operators (<, >, ==)",
        "Combine sensor reading with a while loop",
        "Program the robot to drive until a condition is met",
    ]
    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(18), color=DARK_GRAY)

    # Agenda
    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    agenda = [
        ("while loops explained", "10 min"),
        ("Guided practice: drive to edge", "15 min"),
        ("Independent practice", "20 min"),
    ]
    for item, duration in agenda:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: The Problem
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Problem")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Last lesson:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "We read sensor values, but the robot just sat still.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Today's goal:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Make the robot drive forward AND check the sensor continuously.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "When it detects the line \u2014 stop!", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, 'Question:', level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, '"How do we repeat \u2018drive and check\u2019 without knowing how many times?"',
                    level=0, font_size=Pt(20), color=ACCENT_BLUE, bold=True)

    # ════════════════════════════════════════
    # SLIDE 3: for Loop vs. while Loop
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "for Loop vs. while Loop")

    # Left column: for loop
    add_label(slide, "for loop (Module 1):  Repeat a FIXED number of times",
              Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5), font_size=Pt(18))

    add_code_block(slide,
        "for i in range(4):    # Always runs 4 times\n"
        "    drivetrain.straight(30)\n"
        "    drivetrain.turn(90)",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.7)
    )

    # Right column: while loop
    add_label(slide, "while loop (NEW):  Repeat UNTIL a condition changes",
              Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5), font_size=Pt(18))

    add_code_block(slide,
        "while left < 0.5:    # Runs until left >= 0.5\n"
        "    drivetrain.set_effort(0.3, 0.3)\n"
        "    left = reflectance.get_left()",
        Inches(7.0), Inches(2.0), Inches(6.0), Inches(1.7)
    )

    # Key difference callout
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key difference:  for = fixed count.     while = until something happens."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 4: while Loop Syntax
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "while Loop Syntax")

    add_code_block(slide,
        "while condition:\n"
        "    # code that runs while condition is True\n"
        "    # code stops when condition becomes False",
        Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.0)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(12), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Rules:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    rules = [
        "Colon : after the condition",
        "Indented body (same as for)",
        "Condition is checked BEFORE each iteration",
        "Loop exits when condition is False",
    ]
    for rule in rules:
        add_bullet_text(tf, rule, level=0, font_size=Pt(20), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 5: Comparison Operators
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Comparison Operators")

    headers = ["Operator", "Meaning", "Example", "Result"]
    rows = [
        ["`<`",  "Less than",         "0.2 < 0.5",  "True"],
        ["`>`",  "Greater than",      "0.8 > 0.5",  "True"],
        ["`==`", "Equal to",          "4 == 4",     "True"],
        ["`!=`", "Not equal",         "3 != 5",     "True"],
        ["`<=`", "Less or equal",     "5 <= 5",     "True"],
        ["`>=`", "Greater or equal",  "0.6 >= 0.5", "True"],
    ]
    add_table(slide, headers, rows, Inches(1.0), Inches(1.5), Inches(11), row_height=Inches(0.42))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "For sensors:  left < 0.5  means \u201csensor is on white (below threshold)\u201d"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 6: Drive to the Edge
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Drive to the Edge")

    add_code_block(slide,
        "threshold = 0.5\n"
        "\n"
        "board.wait_for_button()\n"
        "left = reflectance.get_left()\n"
        "\n"
        "while left < threshold:\n"
        "    drivetrain.set_effort(0.3, 0.3)\n"
        "    left = reflectance.get_left()  # UPDATE!\n"
        "\n"
        "drivetrain.stop()\n"
        'print("Line detected!")',
        Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.5)
    )

    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(6.0), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Step by step:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    steps = [
        "Read sensor before entering loop",
        "While on white (below threshold) \u2192 drive forward",
        "Update sensor reading EVERY iteration",
        "When sensor hits line (above threshold) \u2192 loop exits \u2192 stop",
    ]
    for i, step in enumerate(steps, 1):
        add_bullet_text(tf, f"{i}. {step}", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 7: CRITICAL — Update Inside the Loop!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "CRITICAL \u2014 Update Inside the Loop!")

    # CORRECT box
    add_label(slide, "CORRECT:", Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
              font_size=Pt(20), color=GREEN)

    add_code_block(slide,
        "while left < threshold:\n"
        "    drivetrain.set_effort(0.3, 0.3)\n"
        "    left = reflectance.get_left()  # Updates!",
        Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.7)
    )

    # WRONG box
    add_label(slide, "WRONG:", Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5),
              font_size=Pt(20), color=RED)

    add_code_block(slide,
        "while left < threshold:\n"
        "    drivetrain.set_effort(0.3, 0.3)\n"
        "    # left never changes \u2192 infinite loop!",
        Inches(7.0), Inches(2.0), Inches(6.0), Inches(1.7)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Rule: Always update the variable you\u2019re checking inside the while loop!"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 8: set_effort vs. straight
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "set_effort vs. straight")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "drivetrain.straight(30)", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Drives exactly 30 cm, then stops.  BLOCKS until done.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "drivetrain.set_effort(0.3, 0.3)", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Sets motor power and returns IMMEDIATELY.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Why set_effort?", level=0, font_size=Pt(22), bold=True, color=ACCENT_BLUE)
    add_bullet_text(tf, "We need to check the sensor continuously.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf,
        "straight() would drive the whole distance before we could check!",
        level=0, font_size=Pt(20), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 9: Your Turn!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Your Turn!")

    # Activity column
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Exercise:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Place robot inside the taped circle.", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Program it to:", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Wait for button press", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Drive forward", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Stop when it detects the line", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Starter code provided \u2014 fill in the missing parts.", level=0, font_size=Pt(18), color=ACCENT_BLUE, bold=True)

    # Checkpoints column
    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Checkpoints:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drives forward?", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Robot stops at the line?", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Print shows sensor values?", level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 10: Connection to Next Lesson
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Next Lesson")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "Today:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drives to the edge and STOPS.", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next lesson (Lesson 3):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drives to the edge and TURNS AROUND!", level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "if/else statements \u2014 making decisions", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "\u201cBounce driving\u201d \u2014 robot bounces between edges of the circle", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Preview:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)

    add_code_block(slide,
        "if left > threshold:\n"
        "    drivetrain.turn(180)  # Turn around!\n"
        "else:\n"
        "    drivetrain.set_effort(0.3, 0.3)  # Keep going",
        Inches(0.6), Inches(5.5), Inches(7.5), Inches(1.7)
    )

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "02-drive-to-the-edge.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
