#!/usr/bin/env python3
"""Generate grid diagrams for course slides.

This script creates grid path visualizations as PNG images
and can insert them into existing PPTX files.
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import numpy as np


def draw_grid_path(filename, rows, cols, path, blocked=None,
                   start_label="START", end_label="END",
                   title=None, figsize=(8, 6)):
    """
    Draw a grid with a highlighted path.

    Parameters:
        filename: output PNG path
        rows, cols: grid dimensions
        path: list of (row, col) tuples
        blocked: list of (row, col) tuples for blocked nodes (optional)
        start_label, end_label: labels for first/last path nodes
        title: optional title text
        figsize: figure size in inches
    """
    if blocked is None:
        blocked = []

    fig, ax = plt.subplots(1, 1, figsize=figsize)

    # Colors
    GRID_COLOR = '#CCCCCC'
    PATH_COLOR = '#2E75B6'
    PATH_FILL = '#D6E8F7'
    BLOCKED_COLOR = '#FF6B6B'
    BLOCKED_FILL = '#FFD4D4'
    NODE_COLOR = '#333333'
    START_COLOR = '#27AE60'
    END_COLOR = '#E74C3C'
    BG_COLOR = '#FAFAFA'

    ax.set_facecolor(BG_COLOR)
    fig.patch.set_facecolor('white')

    # Draw grid lines
    for r in range(rows):
        for c in range(cols):
            # Horizontal edges
            if c < cols - 1:
                ax.plot([c, c + 1], [rows - 1 - r, rows - 1 - r],
                        color=GRID_COLOR, linewidth=1.5, zorder=1)
            # Vertical edges
            if r < rows - 1:
                ax.plot([c, c], [rows - 1 - r, rows - 2 - r],
                        color=GRID_COLOR, linewidth=1.5, zorder=1)

    # Draw blocked nodes
    for (br, bc) in blocked:
        y = rows - 1 - br
        circle = plt.Circle((bc, y), 0.25, facecolor=BLOCKED_FILL,
                             edgecolor=BLOCKED_COLOR, linewidth=2, zorder=3)
        ax.add_patch(circle)
        ax.text(bc, y, 'X', ha='center', va='center',
                fontsize=14, fontweight='bold', color=BLOCKED_COLOR, zorder=4)

    # Draw path edges (thick colored lines with arrows)
    for i in range(len(path) - 1):
        r1, c1 = path[i]
        r2, c2 = path[i + 1]
        y1 = rows - 1 - r1
        y2 = rows - 1 - r2
        # Draw thick path line
        ax.annotate('', xy=(c2, y2), xytext=(c1, y1),
                    arrowprops=dict(arrowstyle='->', color=PATH_COLOR,
                                    lw=3, mutation_scale=20),
                    zorder=5)

    # Draw all grid nodes
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue

            is_on_path = (r, c) in path
            is_start = path and (r, c) == path[0]
            is_end = path and (r, c) == path[-1]

            if is_start:
                color = START_COLOR
                size = 0.28
            elif is_end:
                color = END_COLOR
                size = 0.28
            elif is_on_path:
                color = PATH_COLOR
                size = 0.22
            else:
                color = NODE_COLOR
                size = 0.12

            if is_on_path:
                circle = plt.Circle((c, y), size, facecolor='white',
                                     edgecolor=color, linewidth=2.5, zorder=6)
                ax.add_patch(circle)
                circle2 = plt.Circle((c, y), size * 0.5, facecolor=color,
                                      edgecolor=color, linewidth=0, zorder=7)
                ax.add_patch(circle2)
            else:
                circle = plt.Circle((c, y), size, facecolor=NODE_COLOR,
                                     edgecolor=NODE_COLOR, linewidth=0, zorder=6)
                ax.add_patch(circle)

    # Add coordinate labels
    for c in range(cols):
        ax.text(c, rows - 0.3, f'Col {c}', ha='center', va='center',
                fontsize=10, color='#666666', fontweight='bold')
    for r in range(rows):
        y = rows - 1 - r
        ax.text(-0.7, y, f'Row {r}', ha='center', va='center',
                fontsize=10, color='#666666', fontweight='bold')

    # Add start/end labels
    if path:
        sr, sc = path[0]
        sy = rows - 1 - sr
        # Place start label to the right if at col 0, otherwise above
        if sc == 0:
            ax.text(sc + 0.5, sy + 0.35, start_label, ha='left', va='bottom',
                    fontsize=11, fontweight='bold', color=START_COLOR, zorder=8)
        else:
            ax.text(sc, sy + 0.5, start_label, ha='center', va='bottom',
                    fontsize=11, fontweight='bold', color=START_COLOR, zorder=8)

        er, ec = path[-1]
        ey = rows - 1 - er
        ax.text(ec, ey - 0.5, end_label, ha='center', va='top',
                fontsize=11, fontweight='bold', color=END_COLOR, zorder=8)

    # Add title
    if title:
        ax.set_title(title, fontsize=14, fontweight='bold', pad=20, color='#1B3A5C')

    # Formatting
    ax.set_xlim(-1, cols)
    ax.set_ylim(-1, rows)
    ax.set_aspect('equal')
    ax.axis('off')

    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def insert_image_into_pptx(pptx_path, slide_number, image_path,
                            left_inches=1.5, top_inches=1.5,
                            width_inches=5, output_path=None):
    """
    Insert an image into an existing PPTX file on a specific slide.

    Parameters:
        pptx_path: path to existing PPTX file
        slide_number: 1-indexed slide number
        image_path: path to PNG image
        left_inches, top_inches: position from top-left
        width_inches: image width (height auto-scaled)
        output_path: where to save (defaults to overwriting pptx_path)
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(pptx_path)
    slide = prs.slides[slide_number - 1]  # Convert to 0-indexed

    slide.shapes.add_picture(
        image_path,
        Inches(left_inches),
        Inches(top_inches),
        Inches(width_inches)
    )

    save_path = output_path or pptx_path
    prs.save(save_path)
    print(f"Inserted {image_path} into slide {slide_number} of {save_path}")


def draw_labeled_grid(filename, rows, cols, blocked=None, title=None,
                      figsize=(8, 6), show_coords=True):
    """
    Draw a grid with coordinate labels at every intersection.

    Parameters:
        filename: output PNG path
        rows, cols: grid dimensions
        blocked: list of (row, col) for blocked nodes
        title: optional title
        figsize: figure size
        show_coords: if True, label each node with (r,c)
    """
    if blocked is None:
        blocked = []

    fig, ax = plt.subplots(1, 1, figsize=figsize)

    GRID_COLOR = '#CCCCCC'
    BLOCKED_COLOR = '#FF6B6B'
    BLOCKED_FILL = '#FFD4D4'
    NODE_COLOR = '#333333'
    BG_COLOR = '#FAFAFA'

    ax.set_facecolor(BG_COLOR)
    fig.patch.set_facecolor('white')

    # Draw grid edges (skip edges to/from blocked nodes)
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue
            if c < cols - 1 and (r, c + 1) not in blocked:
                ax.plot([c, c + 1], [y, y], color=GRID_COLOR, linewidth=1.5, zorder=1)
            if r < rows - 1 and (r + 1, c) not in blocked:
                ax.plot([c, c], [y, y - 1], color=GRID_COLOR, linewidth=1.5, zorder=1)

    # Draw blocked nodes
    for (br, bc) in blocked:
        y = rows - 1 - br
        circle = plt.Circle((bc, y), 0.25, facecolor=BLOCKED_FILL,
                             edgecolor=BLOCKED_COLOR, linewidth=2, zorder=3)
        ax.add_patch(circle)
        ax.text(bc, y, 'X', ha='center', va='center',
                fontsize=14, fontweight='bold', color=BLOCKED_COLOR, zorder=4)

    # Draw nodes with coordinate labels
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue
            circle = plt.Circle((c, y), 0.22, facecolor='white',
                                 edgecolor=NODE_COLOR, linewidth=1.5, zorder=5)
            ax.add_patch(circle)
            if show_coords:
                ax.text(c, y, f'({r},{c})', ha='center', va='center',
                        fontsize=7, color=NODE_COLOR, fontweight='bold', zorder=6)

    # Row/col headers
    for c in range(cols):
        ax.text(c, rows - 0.3, f'Col {c}', ha='center', va='center',
                fontsize=10, color='#666666', fontweight='bold')
    for r in range(rows):
        y = rows - 1 - r
        ax.text(-0.7, y, f'Row {r}', ha='center', va='center',
                fontsize=10, color='#666666', fontweight='bold')

    if title:
        ax.set_title(title, fontsize=14, fontweight='bold', pad=20, color='#1B3A5C')

    ax.set_xlim(-1, cols)
    ax.set_ylim(-0.7, rows)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_side_by_side_grids(filename, rows, cols,
                            path_left, path_right,
                            blocked_left=None, blocked_right=None,
                            title_left="Without Obstacles",
                            title_right="With Obstacle",
                            suptitle=None, figsize=(14, 5)):
    """
    Draw two grids side by side for comparison.
    """
    if blocked_left is None:
        blocked_left = []
    if blocked_right is None:
        blocked_right = []

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=figsize)
    fig.patch.set_facecolor('white')

    for ax, path, blocked, subtitle in [
        (ax1, path_left, blocked_left, title_left),
        (ax2, path_right, blocked_right, title_right)
    ]:
        GRID_COLOR = '#CCCCCC'
        PATH_COLOR = '#2E75B6'
        BLOCKED_COLOR = '#FF6B6B'
        BLOCKED_FILL = '#FFD4D4'
        NODE_COLOR = '#333333'
        START_COLOR = '#27AE60'
        END_COLOR = '#E74C3C'
        BG_COLOR = '#FAFAFA'

        ax.set_facecolor(BG_COLOR)

        # Grid lines
        for r in range(rows):
            for c in range(cols):
                y = rows - 1 - r
                if c < cols - 1:
                    ax.plot([c, c + 1], [y, y], color=GRID_COLOR, linewidth=1.5, zorder=1)
                if r < rows - 1:
                    ax.plot([c, c], [y, y - 1], color=GRID_COLOR, linewidth=1.5, zorder=1)

        # Blocked
        for (br, bc) in blocked:
            y = rows - 1 - br
            circle = plt.Circle((bc, y), 0.25, facecolor=BLOCKED_FILL,
                                 edgecolor=BLOCKED_COLOR, linewidth=2, zorder=3)
            ax.add_patch(circle)
            ax.text(bc, y, 'X', ha='center', va='center',
                    fontsize=14, fontweight='bold', color=BLOCKED_COLOR, zorder=4)

        # Path arrows
        for i in range(len(path) - 1):
            r1, c1 = path[i]
            r2, c2 = path[i + 1]
            y1 = rows - 1 - r1
            y2 = rows - 1 - r2
            ax.annotate('', xy=(c2, y2), xytext=(c1, y1),
                        arrowprops=dict(arrowstyle='->', color=PATH_COLOR,
                                        lw=3, mutation_scale=20), zorder=5)

        # Nodes
        for r in range(rows):
            for c in range(cols):
                y = rows - 1 - r
                if (r, c) in blocked:
                    continue
                is_on_path = (r, c) in path
                is_start = path and (r, c) == path[0]
                is_end = path and (r, c) == path[-1]

                if is_start:
                    color = START_COLOR
                    size = 0.25
                elif is_end:
                    color = END_COLOR
                    size = 0.25
                elif is_on_path:
                    color = PATH_COLOR
                    size = 0.2
                else:
                    color = NODE_COLOR
                    size = 0.1

                if is_on_path:
                    circle = plt.Circle((c, y), size, facecolor='white',
                                         edgecolor=color, linewidth=2.5, zorder=6)
                    ax.add_patch(circle)
                    circle2 = plt.Circle((c, y), size * 0.5, facecolor=color,
                                          edgecolor=color, linewidth=0, zorder=7)
                    ax.add_patch(circle2)
                else:
                    circle = plt.Circle((c, y), size, facecolor=NODE_COLOR,
                                         edgecolor=NODE_COLOR, linewidth=0, zorder=6)
                    ax.add_patch(circle)

        ax.set_title(subtitle, fontsize=12, fontweight='bold', color='#1B3A5C', pad=10)
        ax.set_xlim(-0.5, cols - 0.5)
        ax.set_ylim(-0.5, rows - 0.5)
        ax.set_aspect('equal')
        ax.axis('off')

    if suptitle:
        fig.suptitle(suptitle, fontsize=14, fontweight='bold', color='#1B3A5C', y=1.02)

    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_dijkstra_step(filename, rows, cols, blocked, distances,
                       visited, current_node, title=None, figsize=(7, 7)):
    """
    Draw a single step of Dijkstra's algorithm showing distances and visited status.
    """
    fig, ax = plt.subplots(1, 1, figsize=figsize)

    GRID_COLOR = '#CCCCCC'
    BLOCKED_COLOR = '#FF6B6B'
    BLOCKED_FILL = '#FFD4D4'
    VISITED_COLOR = '#27AE60'
    VISITED_FILL = '#D5F5E3'
    CURRENT_COLOR = '#F39C12'
    CURRENT_FILL = '#FEF9E7'
    KNOWN_COLOR = '#2E75B6'
    KNOWN_FILL = '#D6E8F7'
    UNKNOWN_COLOR = '#999999'
    BG_COLOR = '#FAFAFA'

    ax.set_facecolor(BG_COLOR)
    fig.patch.set_facecolor('white')

    # Grid edges (skip blocked)
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue
            if c < cols - 1 and (r, c + 1) not in blocked:
                ax.plot([c, c + 1], [y, y], color=GRID_COLOR, linewidth=1.5, zorder=1)
            if r < rows - 1 and (r + 1, c) not in blocked:
                ax.plot([c, c], [y, y - 1], color=GRID_COLOR, linewidth=1.5, zorder=1)

    # Blocked
    for (br, bc) in blocked:
        y = rows - 1 - br
        circle = plt.Circle((bc, y), 0.3, facecolor=BLOCKED_FILL,
                             edgecolor=BLOCKED_COLOR, linewidth=2, zorder=3)
        ax.add_patch(circle)
        ax.text(bc, y, 'X', ha='center', va='center',
                fontsize=14, fontweight='bold', color=BLOCKED_COLOR, zorder=4)

    # Draw nodes with distance labels
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue

            node = (r, c)
            if node == current_node:
                fill = CURRENT_FILL
                edge = CURRENT_COLOR
            elif node in visited:
                fill = VISITED_FILL
                edge = VISITED_COLOR
            elif node in distances:
                fill = KNOWN_FILL
                edge = KNOWN_COLOR
            else:
                fill = '#F0F0F0'
                edge = UNKNOWN_COLOR

            circle = plt.Circle((c, y), 0.32, facecolor=fill,
                                 edgecolor=edge, linewidth=2.5, zorder=5)
            ax.add_patch(circle)

            if node in distances:
                ax.text(c, y, str(distances[node]), ha='center', va='center',
                        fontsize=14, fontweight='bold', color=edge, zorder=6)
            else:
                ax.text(c, y, '∞', ha='center', va='center',
                        fontsize=14, color=UNKNOWN_COLOR, zorder=6)

    # Coordinate labels
    for c in range(cols):
        ax.text(c, rows - 0.3, f'Col {c}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')
    for r in range(rows):
        y = rows - 1 - r
        ax.text(-0.7, y, f'Row {r}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')

    # Legend
    legend_items = [
        (CURRENT_FILL, CURRENT_COLOR, 'Current'),
        (VISITED_FILL, VISITED_COLOR, 'Visited'),
        (KNOWN_FILL, KNOWN_COLOR, 'Known dist.'),
        ('#F0F0F0', UNKNOWN_COLOR, 'Unknown'),
    ]
    for i, (fill, edge, label) in enumerate(legend_items):
        lx = cols + 0.3
        ly = rows - 1 - i * 0.6
        circle = plt.Circle((lx, ly), 0.15, facecolor=fill,
                             edgecolor=edge, linewidth=2, zorder=5)
        ax.add_patch(circle)
        ax.text(lx + 0.35, ly, label, ha='left', va='center',
                fontsize=10, color='#333333', zorder=6)

    if title:
        ax.set_title(title, fontsize=13, fontweight='bold', pad=15, color='#1B3A5C')

    ax.set_xlim(-1, cols + 2)
    ax.set_ylim(-0.7, rows)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_path_with_heading(filename, rows, cols, path, headings,
                           blocked=None, title=None, figsize=(8, 6)):
    """
    Draw a path showing robot heading at each step with direction arrows.
    """
    if blocked is None:
        blocked = []

    fig, ax = plt.subplots(1, 1, figsize=figsize)

    GRID_COLOR = '#CCCCCC'
    PATH_COLOR = '#2E75B6'
    NODE_COLOR = '#333333'
    HEADING_COLOR = '#E67E22'
    BG_COLOR = '#FAFAFA'

    ax.set_facecolor(BG_COLOR)
    fig.patch.set_facecolor('white')

    # Grid lines
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if c < cols - 1:
                ax.plot([c, c + 1], [y, y], color=GRID_COLOR, linewidth=1.5, zorder=1)
            if r < rows - 1:
                ax.plot([c, c], [y, y - 1], color=GRID_COLOR, linewidth=1.5, zorder=1)

    # Path arrows
    for i in range(len(path) - 1):
        r1, c1 = path[i]
        r2, c2 = path[i + 1]
        y1 = rows - 1 - r1
        y2 = rows - 1 - r2
        ax.annotate('', xy=(c2, y2), xytext=(c1, y1),
                    arrowprops=dict(arrowstyle='->', color=PATH_COLOR,
                                    lw=3, mutation_scale=20), zorder=5)

    # Heading arrows (showing which way robot faces)
    heading_dx = {'N': 0, 'S': 0, 'E': 0.3, 'W': -0.3}
    heading_dy = {'N': 0.3, 'S': -0.3, 'E': 0, 'W': 0}

    for i, ((r, c), heading) in enumerate(zip(path, headings)):
        y = rows - 1 - r
        is_start = (i == 0)
        is_end = (i == len(path) - 1)

        color = '#27AE60' if is_start else ('#E74C3C' if is_end else PATH_COLOR)
        size = 0.25

        circle = plt.Circle((c, y), size, facecolor='white',
                             edgecolor=color, linewidth=2.5, zorder=6)
        ax.add_patch(circle)

        # Step number
        ax.text(c, y, str(i + 1), ha='center', va='center',
                fontsize=10, fontweight='bold', color=color, zorder=7)

        # Heading indicator
        dx = heading_dx[heading]
        dy = heading_dy[heading]
        ax.annotate('', xy=(c + dx, y + dy), xytext=(c, y),
                    arrowprops=dict(arrowstyle='->', color=HEADING_COLOR,
                                    lw=2, mutation_scale=15), zorder=8)

    # Non-path nodes
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) not in path:
                circle = plt.Circle((c, y), 0.1, facecolor=NODE_COLOR,
                                     edgecolor=NODE_COLOR, linewidth=0, zorder=4)
                ax.add_patch(circle)

    # Labels
    for c in range(cols):
        ax.text(c, rows - 0.3, f'Col {c}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')
    for r in range(rows):
        y = rows - 1 - r
        ax.text(-0.7, y, f'Row {r}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')

    # Legend for heading arrows
    ax.text(cols + 0.2, rows - 1, '→ = heading', ha='left', va='center',
            fontsize=10, color=HEADING_COLOR, fontweight='bold')

    if title:
        ax.set_title(title, fontsize=13, fontweight='bold', pad=15, color='#1B3A5C')

    ax.set_xlim(-1, cols + 2)
    ax.set_ylim(-0.7, rows)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_compass_rose(filename, title=None, figsize=(6, 6), show_turn_table=False):
    """
    Draw a compass rose showing N/S/E/W with optional turn table.
    """
    fig, ax = plt.subplots(1, 1, figsize=figsize)
    fig.patch.set_facecolor('white')
    ax.set_facecolor('#FAFAFA')

    # Compass circle
    circle = plt.Circle((0, 0), 1.5, facecolor='#F0F4F8', edgecolor='#2E75B6',
                         linewidth=3, zorder=1)
    ax.add_patch(circle)

    # Direction arrows
    dirs = {
        'N': (0, 1.2, '↑ North\nrow - 1'),
        'S': (0, -1.2, '↓ South\nrow + 1'),
        'E': (1.2, 0, 'East →\ncol + 1'),
        'W': (-1.2, 0, '← West\ncol - 1'),
    }
    colors = {'N': '#2E75B6', 'S': '#E74C3C', 'E': '#27AE60', 'W': '#F39C12'}

    for d, (dx, dy, label) in dirs.items():
        ax.annotate('', xy=(dx, dy), xytext=(0, 0),
                    arrowprops=dict(arrowstyle='->', color=colors[d],
                                    lw=4, mutation_scale=25), zorder=3)
        # Place labels outside the circle
        label_x = dx * 1.5
        label_y = dy * 1.5
        ax.text(label_x, label_y, label, ha='center', va='center',
                fontsize=11, fontweight='bold', color=colors[d], zorder=4)

    # Center dot
    ax.plot(0, 0, 'o', markersize=10, color='#333333', zorder=5)

    if title:
        ax.set_title(title, fontsize=14, fontweight='bold', pad=20, color='#1B3A5C')

    ax.set_xlim(-2.5, 2.5)
    ax.set_ylim(-2.5, 2.5)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_turn_table(filename, title=None, figsize=(8, 6)):
    """
    Draw the turn lookup table as a visual grid with colored cells.
    """
    fig, ax = plt.subplots(1, 1, figsize=figsize)
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')

    headings = ['N', 'S', 'E', 'W']
    turns = {
        ('N', 'N'): ('None', '#D5F5E3'),
        ('N', 'S'): ('180°', '#FADBD8'),
        ('N', 'E'): ('Right', '#D6EAF8'),
        ('N', 'W'): ('Left', '#FEF9E7'),
        ('S', 'N'): ('180°', '#FADBD8'),
        ('S', 'S'): ('None', '#D5F5E3'),
        ('S', 'E'): ('Left', '#FEF9E7'),
        ('S', 'W'): ('Right', '#D6EAF8'),
        ('E', 'N'): ('Left', '#FEF9E7'),
        ('E', 'S'): ('Right', '#D6EAF8'),
        ('E', 'E'): ('None', '#D5F5E3'),
        ('E', 'W'): ('180°', '#FADBD8'),
        ('W', 'N'): ('Right', '#D6EAF8'),
        ('W', 'S'): ('Left', '#FEF9E7'),
        ('W', 'E'): ('180°', '#FADBD8'),
        ('W', 'W'): ('None', '#D5F5E3'),
    }

    cell_size = 1.2
    # Column headers
    for j, h in enumerate(headings):
        ax.text((j + 1) * cell_size + cell_size / 2, 0.6,
                f'Need {h}', ha='center', va='center',
                fontsize=11, fontweight='bold', color='#1B3A5C')
    # Row headers
    for i, h in enumerate(headings):
        ax.text(cell_size / 2, -(i + 0.5) * cell_size,
                f'Facing {h}', ha='center', va='center',
                fontsize=11, fontweight='bold', color='#1B3A5C')

    # Cells
    for i, current in enumerate(headings):
        for j, need in enumerate(headings):
            x = (j + 1) * cell_size
            y = -(i) * cell_size - cell_size
            turn_text, color = turns[(current, need)]

            rect = mpatches.FancyBboxPatch(
                (x + 0.05, y + 0.05), cell_size - 0.1, cell_size - 0.1,
                boxstyle="round,pad=0.05", facecolor=color,
                edgecolor='#CCCCCC', linewidth=1.5)
            ax.add_patch(rect)
            ax.text(x + cell_size / 2, y + cell_size / 2, turn_text,
                    ha='center', va='center', fontsize=12,
                    fontweight='bold', color='#333333')

    if title:
        ax.set_title(title, fontsize=14, fontweight='bold', pad=20, color='#1B3A5C')

    ax.set_xlim(-0.2, (len(headings) + 1) * cell_size + 0.2)
    ax.set_ylim(-len(headings) * cell_size - 0.3, 1.2)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


def draw_backpointer_path(filename, rows, cols, blocked, previous,
                          start, end, title=None, figsize=(7, 7)):
    """
    Draw the path reconstruction from Dijkstra's previous dictionary.
    Shows backpointer arrows and the final highlighted path.
    """
    fig, ax = plt.subplots(1, 1, figsize=figsize)

    PATH_COLOR = '#2E75B6'
    BACK_COLOR = '#CCCCCC'
    BLOCKED_COLOR = '#FF6B6B'
    BLOCKED_FILL = '#FFD4D4'
    NODE_COLOR = '#333333'
    START_COLOR = '#27AE60'
    END_COLOR = '#E74C3C'
    BG_COLOR = '#FAFAFA'

    ax.set_facecolor(BG_COLOR)
    fig.patch.set_facecolor('white')

    # Grid edges (skip blocked)
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue
            if c < cols - 1 and (r, c + 1) not in blocked:
                ax.plot([c, c + 1], [y, y], color='#E8E8E8', linewidth=1, zorder=1)
            if r < rows - 1 and (r + 1, c) not in blocked:
                ax.plot([c, c], [y, y - 1], color='#E8E8E8', linewidth=1, zorder=1)

    # Blocked
    for (br, bc) in blocked:
        y = rows - 1 - br
        circle = plt.Circle((bc, y), 0.25, facecolor=BLOCKED_FILL,
                             edgecolor=BLOCKED_COLOR, linewidth=2, zorder=3)
        ax.add_patch(circle)
        ax.text(bc, y, 'X', ha='center', va='center',
                fontsize=14, fontweight='bold', color=BLOCKED_COLOR, zorder=4)

    # Draw all backpointer arrows (faded)
    for node, prev in previous.items():
        if prev is not None:
            r1, c1 = node
            r2, c2 = prev
            y1 = rows - 1 - r1
            y2 = rows - 1 - r2
            ax.annotate('', xy=(c2, y2), xytext=(c1, y1),
                        arrowprops=dict(arrowstyle='->', color=BACK_COLOR,
                                        lw=1.5, mutation_scale=15), zorder=2)

    # Reconstruct and highlight the actual path
    path = []
    current = end
    while current is not None:
        path.append(current)
        current = previous.get(current)
    path.reverse()

    # Draw path arrows (bold)
    for i in range(len(path) - 1):
        r1, c1 = path[i]
        r2, c2 = path[i + 1]
        y1 = rows - 1 - r1
        y2 = rows - 1 - r2
        ax.annotate('', xy=(c2, y2), xytext=(c1, y1),
                    arrowprops=dict(arrowstyle='->', color=PATH_COLOR,
                                    lw=3.5, mutation_scale=22), zorder=5)

    # Draw nodes
    for r in range(rows):
        for c in range(cols):
            y = rows - 1 - r
            if (r, c) in blocked:
                continue
            is_on_path = (r, c) in path
            is_start = (r, c) == start
            is_end_node = (r, c) == end

            if is_start:
                color = START_COLOR
                size = 0.25
            elif is_end_node:
                color = END_COLOR
                size = 0.25
            elif is_on_path:
                color = PATH_COLOR
                size = 0.2
            else:
                color = NODE_COLOR
                size = 0.12

            if is_on_path:
                circle = plt.Circle((c, y), size, facecolor='white',
                                     edgecolor=color, linewidth=2.5, zorder=6)
                ax.add_patch(circle)
                circle2 = plt.Circle((c, y), size * 0.5, facecolor=color,
                                      edgecolor=color, linewidth=0, zorder=7)
                ax.add_patch(circle2)
            else:
                circle = plt.Circle((c, y), size, facecolor='#DDDDDD',
                                     edgecolor='#BBBBBB', linewidth=1, zorder=4)
                ax.add_patch(circle)

    # Labels
    for c in range(cols):
        ax.text(c, rows - 0.3, f'Col {c}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')
    for r in range(rows):
        y = rows - 1 - r
        ax.text(-0.7, y, f'Row {r}', ha='center', va='center',
                fontsize=9, color='#666666', fontweight='bold')

    # Start/End labels
    sr, sc = start
    sy = rows - 1 - sr
    ax.text(sc + 0.5, sy + 0.35, 'START', ha='left', va='bottom',
            fontsize=11, fontweight='bold', color=START_COLOR, zorder=8)
    er, ec = end
    ey = rows - 1 - er
    ax.text(ec, ey - 0.5, 'END', ha='center', va='top',
            fontsize=11, fontweight='bold', color=END_COLOR, zorder=8)

    if title:
        ax.set_title(title, fontsize=13, fontweight='bold', pad=15, color='#1B3A5C')

    ax.set_xlim(-1, cols)
    ax.set_ylim(-0.7, rows)
    ax.set_aspect('equal')
    ax.axis('off')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Generated: {filename}")


if __name__ == '__main__':
    import os

    # Ensure image directories exist
    for d in [
        'module-04-manhattan/slides/images',
        'module-05-dijkstra/slides/images',
    ]:
        os.makedirs(d, exist_ok=True)

    # =========================================================
    # GRAPHIC 1: Module 4, Lesson 7, Slide 5 — Turn Table
    # =========================================================
    draw_turn_table(
        'module-04-manhattan/slides/images/turn-lookup-table.png',
        title='Turn Lookup Table',
        figsize=(8, 6)
    )

    # =========================================================
    # GRAPHIC 2: Module 4, Lesson 7, Slide 7 — Path with Headings
    # =========================================================
    draw_path_with_heading(
        'module-04-manhattan/slides/images/path-with-headings.png',
        rows=3, cols=3,
        path=[(0, 0), (1, 0), (2, 0), (2, 1), (2, 2)],
        headings=['N', 'S', 'S', 'E', 'E'],  # heading AFTER arriving at each node
        title='Walking Through a Path — Tracking Heading',
        figsize=(8, 6)
    )

    # =========================================================
    # GRAPHIC 3: Module 5, Lesson 1, Slide 3 — Manhattan Fails Side-by-Side
    # =========================================================
    draw_side_by_side_grids(
        'module-05-dijkstra/slides/images/manhattan-fails-with-obstacles.png',
        rows=3, cols=3,
        path_left=[(0, 0), (1, 0), (2, 0), (2, 1), (2, 2)],
        path_right=[(0, 0)],  # Manhattan can't complete — just show start
        blocked_left=[],
        blocked_right=[(1, 0)],
        title_left='Manhattan Path (no obstacles)',
        title_right='Manhattan BLOCKED!',
        suptitle='Why Manhattan Fails with Obstacles',
        figsize=(12, 5)
    )

    # =========================================================
    # GRAPHIC 4: Module 5, Lesson 1, Slide 5 — 3x3 Graph
    # =========================================================
    draw_labeled_grid(
        'module-05-dijkstra/slides/images/grid-as-graph-3x3.png',
        rows=3, cols=3,
        title='The Grid as a Graph — 3×3',
        figsize=(7, 6)
    )

    # =========================================================
    # GRAPHIC 5: Module 5, Lesson 3, Slides 5-7 — Dijkstra Steps
    # =========================================================
    blocked_11 = [(1, 1)]

    # Step 0: Setup
    draw_dijkstra_step(
        'module-05-dijkstra/slides/images/dijkstra-step0-setup.png',
        rows=3, cols=3, blocked=blocked_11,
        distances={(0, 0): 0},
        visited=[],
        current_node=(0, 0),
        title='Dijkstra Setup — Start at (0,0)',
        figsize=(7, 6)
    )

    # Step 1: Visit (0,0)
    draw_dijkstra_step(
        'module-05-dijkstra/slides/images/dijkstra-step1.png',
        rows=3, cols=3, blocked=blocked_11,
        distances={(0, 0): 0, (0, 1): 1, (1, 0): 1},
        visited=[(0, 0)],
        current_node=(0, 1),
        title='Step 1: Visited (0,0) → neighbors get distance 1',
        figsize=(7, 6)
    )

    # Step 2: Visit (0,1)
    draw_dijkstra_step(
        'module-05-dijkstra/slides/images/dijkstra-step2.png',
        rows=3, cols=3, blocked=blocked_11,
        distances={(0, 0): 0, (0, 1): 1, (1, 0): 1, (0, 2): 2},
        visited=[(0, 0), (0, 1)],
        current_node=(1, 0),
        title='Step 2: Visited (0,1) → (0,2) gets distance 2',
        figsize=(7, 6)
    )

    # Step 3: Visit (1,0)
    draw_dijkstra_step(
        'module-05-dijkstra/slides/images/dijkstra-step3.png',
        rows=3, cols=3, blocked=blocked_11,
        distances={(0, 0): 0, (0, 1): 1, (1, 0): 1, (0, 2): 2, (2, 0): 2},
        visited=[(0, 0), (0, 1), (1, 0)],
        current_node=(0, 2),
        title='Step 3: Visited (1,0) → (2,0) gets distance 2',
        figsize=(7, 6)
    )

    # Steps 4-6: Final state
    draw_dijkstra_step(
        'module-05-dijkstra/slides/images/dijkstra-final.png',
        rows=3, cols=3, blocked=blocked_11,
        distances={
            (0, 0): 0, (0, 1): 1, (0, 2): 2,
            (1, 0): 1, (1, 2): 3,
            (2, 0): 2, (2, 1): 3, (2, 2): 4
        },
        visited=[(0, 0), (0, 1), (1, 0), (0, 2), (2, 0), (1, 2), (2, 1), (2, 2)],
        current_node=None,
        title='Final: All distances from (0,0)',
        figsize=(7, 6)
    )

    # =========================================================
    # GRAPHIC 6: Module 5, Lesson 3, Slide 8 — Backpointer Path
    # =========================================================
    previous = {
        (0, 0): None,
        (0, 1): (0, 0),
        (0, 2): (0, 1),
        (1, 0): (0, 0),
        (1, 2): (0, 2),
        (2, 0): (1, 0),
        (2, 1): (2, 0),
        (2, 2): (1, 2),
    }
    draw_backpointer_path(
        'module-05-dijkstra/slides/images/backpointer-path.png',
        rows=3, cols=3, blocked=blocked_11,
        previous=previous,
        start=(0, 0), end=(2, 2),
        title='Reconstructing the Path: (0,0) → (2,2)',
        figsize=(7, 6)
    )

    # =========================================================
    # GRAPHIC 7: Module 4, Lesson 4, Slide 6 — Compass Rose with Directions
    # =========================================================
    draw_compass_rose(
        'module-04-manhattan/slides/images/compass-rose-directions.png',
        title='Grid Directions & Row/Column Changes',
        figsize=(6, 6)
    )

    # =========================================================
    # GRAPHIC 8: Module 4, Lesson 1, Slide 4 — Labeled 4x4 Grid
    # =========================================================
    draw_labeled_grid(
        'module-04-manhattan/slides/images/labeled-grid-4x4.png',
        rows=4, cols=4,
        title='Grid Coordinates — (row, col)',
        figsize=(8, 7)
    )

    # =========================================================
    # GRAPHIC 9: Module 5, Lesson 4, Slide 9 — 4x4 Graph with Blocked Node
    # =========================================================
    draw_labeled_grid(
        'module-05-dijkstra/slides/images/graph-4x4-blocked.png',
        rows=4, cols=4,
        blocked=[(1, 1)],
        title='4×4 Graph with (1,1) Blocked',
        figsize=(8, 7)
    )

    # =========================================================
    # GRAPHIC 10: Module 5, Lesson 9, Slide 6 — Run 1 vs Run 2 Comparison
    # =========================================================
    obstacles = [(1, 1), (2, 3)]
    # Run 1: Discovery run - takes a longer path
    run1_path = [(0, 0), (1, 0), (2, 0), (2, 1), (2, 2), (1, 2), (0, 2), (0, 3), (1, 3), (2, 3)]
    # Oops (2,3) is blocked — reroute
    run1_path = [(0, 0), (1, 0), (2, 0), (2, 1), (2, 2), (1, 2), (1, 3)]
    # Run 2: Knows obstacles — plans optimal path
    run2_path = [(0, 0), (0, 1), (0, 2), (1, 2), (1, 3)]

    draw_side_by_side_grids(
        'module-05-dijkstra/slides/images/run1-vs-run2.png',
        rows=3, cols=4,
        path_left=run1_path,
        path_right=run2_path,
        blocked_left=obstacles,
        blocked_right=obstacles,
        title_left='Run 1: Discovery (7 steps)',
        title_right='Run 2: With Experience (5 steps)',
        suptitle='Learning from Experience',
        figsize=(14, 5)
    )

    # === Original proof of concept diagrams ===
    draw_grid_path(
        'module-04-manhattan/slides/images/path-0_0-to-2_3.png',
        rows=3, cols=4,
        path=[(0, 0), (1, 0), (2, 0), (2, 1), (2, 2), (2, 3)],
        title='Manhattan Path: (0, 0) → (2, 3)',
        figsize=(8, 5)
    )

    draw_grid_path(
        'module-05-dijkstra/slides/images/dijkstra-with-obstacles.png',
        rows=4, cols=4,
        path=[(0, 0), (0, 1), (0, 2), (1, 2), (2, 2), (2, 1), (2, 0), (3, 0)],
        blocked=[(1, 0), (1, 1)],
        title='Dijkstra Path Around Obstacles',
        figsize=(7, 7)
    )
