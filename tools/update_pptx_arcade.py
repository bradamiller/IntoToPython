#!/usr/bin/env python3
"""Update PPTX slides to use arcade() instead of set_effort() for proportional control.

Only modifies the specific slides that contain the old set_effort patterns.
Does NOT regenerate slides — preserves all existing formatting and corrections.
"""

from pptx import Presentation
from pptx.util import Pt
import os
import re
import copy


def replace_in_paragraph(paragraph, old_text, new_text):
    """Replace text in a paragraph while trying to preserve formatting.

    This works by finding the run(s) containing the old text and replacing.
    For simple cases where the text is in a single run, formatting is preserved.
    """
    full_text = ''.join(run.text for run in paragraph.runs)
    if old_text not in full_text:
        return False

    new_full = full_text.replace(old_text, new_text)

    # If we have runs, put all text in the first run and clear others
    if paragraph.runs:
        # Preserve first run's formatting
        paragraph.runs[0].text = new_full
        for run in paragraph.runs[1:]:
            run.text = ''
    return True


def replace_in_shape(shape, old_text, new_text):
    """Replace text in all paragraphs of a shape."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        if replace_in_paragraph(para, old_text, new_text):
            changed = True
    return changed


def replace_in_slide(slide, replacements):
    """Apply a list of (old, new) text replacements to a slide."""
    changed = False
    for shape in slide.shapes:
        for old_text, new_text in replacements:
            if replace_in_shape(shape, old_text, new_text):
                changed = True
    return changed


def update_file(pptx_path, slide_updates):
    """Update specific slides in a PPTX file.

    slide_updates: dict mapping slide_number (1-based) to list of (old, new) replacements
    """
    prs = Presentation(pptx_path)
    any_changed = False

    for slide_num, replacements in slide_updates.items():
        slide = prs.slides[slide_num - 1]
        if replace_in_slide(slide, replacements):
            any_changed = True
            print(f"  Updated slide {slide_num}")
        else:
            # Try checking what text is actually there
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for old_text, _ in replacements:
                        if old_text in shape.text_frame.text:
                            print(f"  WARNING: Found '{old_text[:50]}' but replacement failed on slide {slide_num}")

    if any_changed:
        prs.save(pptx_path)
        print(f"  Saved: {pptx_path}")
    else:
        print(f"  No changes needed: {pptx_path}")


def main():
    base = 'module-02-line-tracking/slides'

    # === 05-proportional-control.pptx ===
    print("05-proportional-control.pptx:")
    update_file(os.path.join(base, '05-proportional-control.pptx'), {
        # Slide 1: Learning objectives and agenda
        1: [
            ('Use set_effort() for continuous motor control',
             'Use arcade() for continuous motor control with steering'),
            ('Introducing set_effort() and continuous control',
             'Introducing arcade() and continuous control'),
        ],
        # Slide 3: Blocking vs continuous
        3: [
            ('drivetrain.set_effort(0.3, 0.3)   # Sets motors and CONTINUES',
             'drivetrain.arcade(0.3, 0)   # Drive forward and KEEP GOING'),
            ('drivetrain.set_effort(0.3, 0.3)   # Sets motors and KEEPS GOING',
             'drivetrain.arcade(0.3, 0)   # Drive forward and KEEP GOING'),
            ('set_effort() does not block',
             'arcade() does not block'),
        ],
        # Slide 4: set_effort examples -> arcade examples
        4: [
            ('set_effort() Examples',
             'arcade() — Speed and Steering'),
            ('drivetrain.set_effort(0.3, 0.3)   # Both motors equal',
             'drivetrain.arcade(0.3, 0)      # Forward, no turning'),
            ('drivetrain.set_effort(0.4, 0.2)   # Left faster than right',
             'drivetrain.arcade(0.3, 0.15)   # Forward + steer right'),
            ('drivetrain.set_effort(0.2, 0.4)   # Right faster than left',
             'drivetrain.arcade(0.3, -0.15)  # Forward + steer left'),
            ('Drive straight:',
             'Drive straight:'),
            ('Turn right (left faster):',
             'Steer right:'),
            ('Turn left (right faster):',
             'Steer left:'),
            ('Values range from -1.0 (full reverse) to 1.0 (full forward)',
             'Under the hood: arcade(speed, turn) sets left motor to speed + turn and right motor to speed - turn.'),
            ('What would set_effort(0.5, -0.5) do?" (Spin in place)',
             'What would arcade(0.3, 0.15) do to each motor?" (Left = 0.45, Right = 0.15)'),
        ],
        # Slide 7: From Error to Correction
        7: [
            ('left_effort  = base_effort + correction\nright_effort = base_effort - correction',
             'drivetrain.arcade(base_effort, correction)'),
            ('Then apply correction to motors:',
             'Then apply correction to motors using arcade:'),
        ],
        # Slide 8: Complete control loop
        8: [
            ('    # 4. Apply to motors\n    left_effort = base_effort + correction\n    right_effort = base_effort - correction\n    drivetrain.set_effort(left_effort, right_effort)',
             '    # 4. Apply to motors\n    drivetrain.arcade(base_effort, correction)'),
        ],
        # Slide 9: Tracing through the math
        9: [
            ('left_effort   = 0.3 + 0.15 = 0.45\nright_effort  = 0.3 - 0.15 = 0.15',
             'arcade(0.3, 0.15) --> left = 0.3 + 0.15 = 0.45\n                      right = 0.3 - 0.15 = 0.15'),
            ('left_effort   = 0.3 + (-0.15) = 0.15\nright_effort  = 0.3 - (-0.15) = 0.45',
             'arcade(0.3, -0.15) --> left = 0.3 + (-0.15) = 0.15\n                       right = 0.3 - (-0.15) = 0.45'),
        ],
        # Slide 11: Exercise code
        11: [
            ('    drivetrain.set_effort(base_effort + correction,\n                          base_effort - correction)',
             '    drivetrain.arcade(base_effort, correction)'),
        ],
    })

    # === 06-two-sensor-line-following.pptx ===
    print("\n06-two-sensor-line-following.pptx:")
    update_file(os.path.join(base, '06-two-sensor-line-following.pptx'), {
        # Slide 5: Applying the correction
        5: [
            ('left_effort  = base_effort - error * Kp\nright_effort = base_effort + error * Kp',
             'correction = error * Kp\ndrivetrain.arcade(base_effort, -correction)'),
            ('Motor effort formula:',
             'Apply the correction with arcade:'),
            ('left gets MINUS correction, right gets PLUS correction',
             'we negate the correction when calling arcade'),
        ],
        # Slide 6: Trace through the math
        6: [
            ('left_effort   = 0.3 - (0.6 * 0.5) = 0.3 - 0.3 = 0.0\nright_effort  = 0.3 + (0.6 * 0.5) = 0.3 + 0.3 = 0.6',
             'correction    = 0.6 * 0.5 = 0.3\narcade(0.3, -0.3) --> left = 0.0, right = 0.6'),
            ('left_effort   = 0.3 - (-0.6 * 0.5) = 0.3 + 0.3 = 0.6\nright_effort  = 0.3 + (-0.6 * 0.5) = 0.3 - 0.3 = 0.0',
             'correction    = -0.6 * 0.5 = -0.3\narcade(0.3, 0.3) --> left = 0.6, right = 0.0'),
        ],
        # Slide 7: Complete program
        7: [
            ('    left_effort = base_effort - error * Kp\n    right_effort = base_effort + error * Kp\n    drivetrain.set_effort(left_effort, right_effort)',
             '    correction = error * Kp\n    drivetrain.arcade(base_effort, -correction)'),
        ],
    })

    # === 07-detecting-intersections.pptx ===
    print("\n07-detecting-intersections.pptx:")
    update_file(os.path.join(base, '07-detecting-intersections.pptx'), {
        5: [
            ('        left_effort = base_effort - error * Kp\n        right_effort = base_effort + error * Kp\n        drivetrain.set_effort(left_effort, right_effort)',
             '        correction = error * Kp\n        drivetrain.arcade(base_effort, -correction)'),
        ],
    })

    # === 08-introduction-to-classes.pptx ===
    print("\n08-introduction-to-classes.pptx:")
    update_file(os.path.join(base, '08-introduction-to-classes.pptx'), {
        2: [
            ('    drivetrain.set_effort(base_effort - error * kp,\n                          base_effort + error * kp)',
             '    correction = error * kp\n    drivetrain.arcade(base_effort, -correction)'),
        ],
    })

    # === 09-object-composition.pptx ===
    print("\n09-object-composition.pptx:")
    update_file(os.path.join(base, '09-object-composition.pptx'), {
        5: [
            ('        left = self.base_effort - error * self.Kp\n        right = self.base_effort + error * self.Kp\n        self.drivetrain.set_effort(left, right)',
             '        correction = error * self.Kp\n        self.drivetrain.arcade(self.base_effort, -correction)'),
        ],
        6: [
            ('    self.drivetrain.set_effort(self.base_effort, self.base_effort)',
             '    self.drivetrain.arcade(self.base_effort, 0)'),
            ('    self.drivetrain.set_effort(0.3, -0.3)',
             '    self.drivetrain.arcade(0, 0.3)'),
        ],
    })


if __name__ == '__main__':
    main()
