#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 6: Two-Sensor Line Following."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(14)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              row_height=Inches(0.45), font_size=Pt(13)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        cw = int(width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = cw

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.0))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 6: Two-Sensor Line Following"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    objectives = [
        "Understand why two sensors are better than one",
        "Calculate error using left and right sensors: error = left \u2212 right",
        "Implement a two-sensor proportional control loop",
        "Compare one-sensor and two-sensor line following",
    ]
    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5.5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    agenda = [
        ("Why one sensor is not enough", "5 min"),
        ("Two-sensor error calculation", "10 min"),
        ("Build and test two-sensor follower", "15 min"),
        ("Compare one-sensor vs. two-sensor", "20 min"),
    ]
    for item, duration in agenda:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(16), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: Hook — The Limitation of One Sensor
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Limitation of One Sensor")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "One sensor (Lesson 5) follows the EDGE:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Problem scenarios:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Robot drifts completely off the line to the right", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Sensor reads low (white)", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "But is the line to the LEFT or to the RIGHT?", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "One sensor cannot tell!", level=1, font_size=Pt(18), bold=True, color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Solution: Use BOTH sensors \u2014 one on each side of the line.",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, '"If you close one eye, can you still judge distances? What about with both eyes?"',
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 3: Two Sensors Straddle the Line
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Two Sensors Straddle the Line")

    headers = ["Position", "Left Sensor", "Right Sensor", "Meaning"]
    rows = [
        ["Centered",      "~0.5 (medium)", "~0.5 (medium)", "Both sensors at edge \u2014 drive straight"],
        ["Drifted left",  "High (0.8+)",   "Low (0.2\u2013)",  "Left on line, right on white \u2014 steer left"],
        ["Drifted right", "Low (0.2\u2013)",  "High (0.8+)",   "Right on line, left on white \u2014 steer right"],
    ]
    col_widths = [Inches(2.2), Inches(2.5), Inches(2.5), Inches(4.5)]
    add_table(slide, headers, rows, Inches(0.4), Inches(1.5), Inches(11.7),
              col_widths=col_widths, row_height=Inches(0.55))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "The difference between sensors tells us EVERYTHING:",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "How far we drifted", level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Which direction we drifted", level=1, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 4: The Two-Sensor Error Formula
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Two-Sensor Error Formula")

    add_label(slide, "One sensor (Lesson 5):", Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    add_code_block(slide,
        "error = setpoint - sensor_reading    # Needs a setpoint!",
        Inches(0.6), Inches(2.0), Inches(12.0), Inches(0.9)
    )

    add_label(slide, "Two sensors (this lesson):", Inches(0.6), Inches(3.0), Inches(12), Inches(0.5))
    add_code_block(slide,
        "error = reflectance.get_left() - reflectance.get_right()    # No setpoint needed!",
        Inches(0.6), Inches(3.6), Inches(12.0), Inches(0.9)
    )

    headers = ["Robot Position", "Left", "Right", "error = L \u2212 R", "Meaning"]
    rows = [
        ["Centered",      "0.5", "0.5", "0.0",  "Drive straight"],
        ["Drifted left",  "0.8", "0.2", "+0.6", "Steer left to correct"],
        ["Drifted right", "0.2", "0.8", "\u22120.6", "Steer right to correct"],
        ["Both on white", "0.1", "0.1", "0.0",  "Off the line"],
        ["Both on black", "0.9", "0.9", "0.0",  "Intersection?"],
    ]
    col_widths = [Inches(2.5), Inches(1.0), Inches(1.0), Inches(1.8), Inches(5.0)]
    add_table(slide, headers, rows, Inches(0.4), Inches(4.7), Inches(11.5),
              col_widths=col_widths, row_height=Inches(0.38))

    # ════════════════════════════════════════
    # SLIDE 5: Applying the Correction
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Applying the Correction")

    add_label(slide, "Motor effort formula:", Inches(0.6), Inches(1.4), Inches(12), Inches(0.5))
    add_code_block(slide,
        "left_effort  = base_effort - error * Kp\n"
        "right_effort = base_effort + error * Kp",
        Inches(0.6), Inches(2.0), Inches(8.0), Inches(1.3)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(12), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Why the signs?", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "If error is positive (drifted left, left sensor on line):",
                    level=0, font_size=Pt(18), bold=True, color=DARK_GRAY)
    add_bullet_text(tf, "Need to steer left to get back", level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Slow down left motor (subtract), speed up right motor (add)",
                    level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "If error is negative (drifted right, right sensor on line):",
                    level=0, font_size=Pt(18), bold=True, color=DARK_GRAY)
    add_bullet_text(tf, "Need to steer right to get back", level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Speed up left motor (subtracting a negative = adding), slow down right motor",
                    level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Note: Signs are reversed from Lesson 5 because the error formula changed.",
                    level=0, font_size=Pt(17), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 6: Trace Through the Math
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Trace Through the Math")

    add_label(slide, "Scenario: Robot drifts left  (left sensor on line)",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5), font_size=Pt(18))
    add_code_block(slide,
        "left_sensor   = 0.8\n"
        "right_sensor  = 0.2\n"
        "error         = 0.8 - 0.2 = 0.6\n"
        "Kp            = 0.5\n"
        "base_effort   = 0.3\n"
        "\n"
        "left_effort   = 0.3 - (0.6 * 0.5) = 0.3 - 0.3 = 0.0\n"
        "right_effort  = 0.3 + (0.6 * 0.5) = 0.3 + 0.3 = 0.6",
        Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.2)
    )

    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(1.9), Inches(4.8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Left motor stops, right runs fast \u2192 robot turns left, back toward center."
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    p.font.name = BODY_FONT

    add_label(slide, "Scenario: Robot drifts right  (right sensor on line)",
              Inches(0.6), Inches(5.3), Inches(12), Inches(0.5), font_size=Pt(18))
    add_code_block(slide,
        "left_sensor   = 0.2\n"
        "right_sensor  = 0.8\n"
        "error         = 0.2 - 0.8 = -0.6\n"
        "\n"
        "left_effort   = 0.3 - (-0.6 * 0.5) = 0.3 + 0.3 = 0.6\n"
        "right_effort  = 0.3 + (-0.6 * 0.5) = 0.3 - 0.3 = 0.0",
        Inches(0.6), Inches(5.9), Inches(7.5), Inches(2.5)
    )

    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(5.9), Inches(4.8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Right motor stops, left runs fast \u2192 robot turns right, back toward center."
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 7: The Complete Two-Sensor Program
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Complete Two-Sensor Program")

    add_code_block(slide,
        "from XRPLib.differential_drive import DifferentialDrive\n"
        "from XRPLib.reflectance import Reflectance\n"
        "from XRPLib.board import Board\n"
        "import time\n"
        "\n"
        "drivetrain  = DifferentialDrive.get_default_differential_drive()\n"
        "reflectance = Reflectance.get_default_reflectance()\n"
        "board       = Board.get_default_board()\n"
        "\n"
        "Kp          = 0.5\n"
        "base_effort = 0.3\n"
        "\n"
        "board.wait_for_button()\n"
        "\n"
        "while True:\n"
        "    # Read both sensors\n"
        "    left_sensor  = reflectance.get_left()\n"
        "    right_sensor = reflectance.get_right()\n"
        "\n"
        "    # Calculate error\n"
        "    error = left_sensor - right_sensor\n"
        "\n"
        "    # Apply correction\n"
        "    left_effort  = base_effort - error * Kp\n"
        "    right_effort = base_effort + error * Kp\n"
        "    drivetrain.set_effort(left_effort, right_effort)\n"
        "\n"
        "    time.sleep(0.01)",
        Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.8),
        font_size=Pt(13)
    )

    # ════════════════════════════════════════
    # SLIDE 8: One Sensor vs. Two Sensors
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "One Sensor vs. Two Sensors")

    headers = ["Feature", "One Sensor (L5)", "Two Sensors (L6)"]
    rows = [
        ["What it follows",  "Edge of the line",               "Center of the line"],
        ["Error formula",    "setpoint \u2212 sensor",              "left \u2212 right"],
        ["Needs setpoint?",  "Yes (0.5)",                      "No"],
        ["Motor formula",    "base \u00b1 correction",               "base \u2213 error\u00d7Kp"],
        ["Smoothness",       "Good",                           "Better"],
        ["Recovery",         "Can only detect one side",       "Detects both sides"],
    ]
    col_widths = [Inches(3.2), Inches(4.0), Inches(4.0)]
    add_table(slide, headers, rows, Inches(0.5), Inches(1.5), Inches(11.5),
              col_widths=col_widths, row_height=Inches(0.5))

    # ════════════════════════════════════════
    # SLIDE 9: Exploring Sensor Readings
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Exploring Sensor Readings")

    add_label(slide, "Before running the control loop, explore the sensor values:",
              Inches(0.6), Inches(1.4), Inches(12), Inches(0.5), font_size=Pt(18))

    add_code_block(slide,
        "from XRPLib.reflectance import Reflectance\n"
        "from XRPLib.board import Board\n"
        "import time\n"
        "\n"
        "reflectance = Reflectance.get_default_reflectance()\n"
        "board = Board.get_default_board()\n"
        "\n"
        "board.wait_for_button()\n"
        "\n"
        "while True:\n"
        "    left  = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "    error = left - right\n"
        '    print(f"L={left:.2f}  R={right:.2f}  err={error:.2f}")\n'
        "    time.sleep(0.2)",
        Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.8)
    )

    txBox = slide.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(4.2), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Activity:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Slowly slide the robot across the line and watch how error changes.",
                    level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Centered: error near 0", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Left on line: error positive", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Right on line: error negative", level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 10: Exercise — Two-Sensor Line Following
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Exercise \u2014 Two-Sensor Line Following")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Your task:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "1. Type in the two-sensor control program", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "2. Place the robot so the line runs between both sensors", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "3. Press button and observe the tracking", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "4. Add debug printing to see the values", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Then compare:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "5. Run your Lesson 5 (one-sensor) program on the same circle", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "6. Run the Lesson 6 (two-sensor) program on the same circle", level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Questions to answer:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Which version follows the line more smoothly?", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Which version handles the curved parts better?", level=0, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "Which version can handle a higher base_effort without losing the line?",
                    level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 11: A Curious Case — Both Sensors High
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "A Curious Case \u2014 Both Sensors High")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "What happens when both sensors read high?", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)

    add_code_block(slide,
        "left_sensor  = 0.9\n"
        "right_sensor = 0.9\n"
        "error        = 0.9 - 0.9 = 0.0",
        Inches(0.6), Inches(2.6), Inches(7.0), Inches(1.8)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Error is zero \u2014 the robot drives straight.",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "But what does it mean physically when both sensors see the line?",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Answer: The robot might be at an INTERSECTION!",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)
    add_bullet_text(tf, "This is the key idea for Lesson 7: both sensors high = intersection detected.",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 12: Connection to Lesson 7
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Lesson 7")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""

    add_bullet_text(tf, "What we accomplished today:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Smooth two-sensor line following", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Error = left \u2212 right (no setpoint needed)", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Smoother tracking than one sensor", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next lesson: Detecting Intersections", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Add a taped cross to the circle", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "When both sensors read high: intersection detected", level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "Robot stops, turns around, and continues in the opposite direction",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, '"How can we tell the difference between \'centered on the line\' and \'at an intersection\'?"',
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    add_code_block(slide,
        "if left_sensor > 0.5 and right_sensor > 0.5:\n"
        "    # We are at an intersection!\n"
        "    drivetrain.stop()\n"
        "    drivetrain.turn(180)",
        Inches(0.6), Inches(5.9), Inches(7.5), Inches(1.5)
    )

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "06-two-sensor-line-following.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
