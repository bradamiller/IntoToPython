#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 8: Introduction to Classes -- LineSensor."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
GREEN = RGBColor(0x17, 0x7B, 0x41)
RED = RGBColor(0xC0, 0x00, 0x00)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(5)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(13)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.1),
        width - Inches(0.4), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              row_height=Inches(0.45), font_size=Pt(13)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        cw = int(width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = cw
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.0))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 8: Introduction to Classes \u2014 LineSensor"
    p.font.size = Pt(34)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for obj in [
        "Understand why classes are useful for organizing code",
        "Write a class with __init__, self, and methods",
        "Create objects and call methods on them",
        "Build a working LineSensor class for the XRP robot",
    ]:
        add_bullet_text(tf, obj, level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5.5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for item, duration in [
        ("Why classes?", "5 min"),
        ("Class syntax walkthrough", "10 min"),
        ("Build LineSensor step by step", "15 min"),
        ("Practice: write and test your LineSensor", "15 min"),
    ]:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(16), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: Hook — The Problem with Our Code
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Problem with Our Code")

    add_label(slide, "Our Lesson 7 code has everything mixed together:",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5), font_size=Pt(18))
    add_code_block(slide,
        "reflectance = Reflectance.get_default_reflectance()\n"
        "drivetrain  = DifferentialDrive.get_default_differential_drive()\n"
        "threshold = 0.5\n"
        "base_effort = 0.4\n"
        "kp = 0.5\n"
        "\n"
        "while True:\n"
        "    left  = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "    error = left - right\n"
        "    if left > threshold and right > threshold:\n"
        "        drivetrain.stop()\n"
        "        break\n"
        "    drivetrain.set_effort(base_effort - error * kp,\n"
        "                          base_effort + error * kp)",
        Inches(0.6), Inches(1.9), Inches(8.5), Inches(4.8),
        font_size=Pt(13)
    )

    txBox = slide.shapes.add_textbox(Inches(9.2), Inches(1.9), Inches(3.9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Problems:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Sensor code and driving code are tangled together",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Hard to reuse just the sensor part in another program",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Hard to read", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, '"How could we separate the sensor logic from the driving logic?"',
                    level=0, font_size=Pt(16), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 3: What Is a Class?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "What Is a Class?")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "A class is a blueprint that groups related data and functions together."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = BODY_FONT

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Analogy \u2014 TV Remote:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Data it stores: current channel, volume level",
                    level=1, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Actions it can do: change channel, adjust volume, power on/off",
                    level=1, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Analogy \u2014 LineSensor:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Data it stores: the reflectance hardware, a threshold value",
                    level=1, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Actions it can do: get error, check for cross, check if off line",
                    level=1, font_size=Pt(16), color=DARK_GRAY)

    headers = ["Term", "Meaning"]
    rows = [
        ["Class",     "The blueprint (the design)"],
        ["Object",    "A specific thing built from the blueprint"],
        ["Method",    "A function that belongs to a class"],
        ["Attribute", "A variable that belongs to an object"],
    ]
    col_widths = [Inches(2.5), Inches(5.5)]
    add_table(slide, headers, rows, Inches(7.0), Inches(2.2), Inches(6.0),
              col_widths=col_widths, row_height=Inches(0.52))

    # ════════════════════════════════════════
    # SLIDE 4: You Already Use Classes!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "You Already Use Classes!")

    add_label(slide, "Every time you wrote this, you were using a class:",
              Inches(0.6), Inches(1.4), Inches(12), Inches(0.5), font_size=Pt(18))

    add_code_block(slide,
        "drivetrain = DifferentialDrive.get_default_differential_drive()\n"
        "drivetrain.straight(30)\n"
        "drivetrain.turn(90)\n"
        "drivetrain.stop()",
        Inches(0.6), Inches(2.0), Inches(10.0), Inches(2.1)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "DifferentialDrive is a class (the blueprint)",
                    level=0, font_size=Pt(19), color=DARK_GRAY)
    add_bullet_text(tf, "drivetrain is an object (a specific instance)",
                    level=0, font_size=Pt(19), color=DARK_GRAY)
    add_bullet_text(tf, ".straight(), .turn(), .stop() are methods",
                    level=0, font_size=Pt(19), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Today you will build your own class \u2014 just like the ones in XRPLib!",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 5: Class Syntax — The Basics
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Class Syntax \u2014 The Basics")

    add_label(slide, "Defining an empty class:", Inches(0.6), Inches(1.4), Inches(6.5), Inches(0.5))
    add_code_block(slide,
        "class LineSensor:\n"
        "    pass",
        Inches(0.6), Inches(2.0), Inches(6.5), Inches(1.1)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(6.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "class \u2014 keyword that starts a class definition",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "LineSensor \u2014 the name (CamelCase: capitalize each word)",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, ": \u2014 colon starts the body (just like def and for)",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "pass \u2014 placeholder for \u201cnothing here yet\u201d",
                    level=0, font_size=Pt(16), color=DARK_GRAY)

    add_label(slide, "Creating an object (instantiation):", Inches(7.2), Inches(1.4), Inches(5.8), Inches(0.5))
    add_code_block(slide,
        "sensor = LineSensor()\n"
        "print(type(sensor))  # <class '__main__.LineSensor'>",
        Inches(7.2), Inches(2.0), Inches(5.8), Inches(1.1)
    )

    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(3.3), Inches(5.8), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key point: Parentheses () after the class name create a new object."
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 6: The __init__ Method (Constructor)
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The __init__ Method (Constructor)")

    add_label(slide, "__init__ runs automatically when you create an object:",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(slide,
        "class LineSensor:\n"
        "    def __init__(self):\n"
        "        self.reflectance = Reflectance.get_default_reflectance()\n"
        "        self.threshold = 0.5",
        Inches(0.6), Inches(1.9), Inches(10.0), Inches(2.0)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.1), Inches(6.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "def __init__(self): \u2014 special method, called the constructor",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "self \u2014 refers to the object being created (\u201cme, myself\u201d)",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "self.reflectance \u2014 stores the sensor hardware on this object",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "self.threshold \u2014 stores 0.5 on this object",
                    level=0, font_size=Pt(16), color=DARK_GRAY)

    add_label(slide, "Using it:", Inches(7.2), Inches(4.1), Inches(5.8), Inches(0.5))
    add_code_block(slide,
        "sensor = LineSensor()       # __init__ runs here!\n"
        "print(sensor.threshold)     # Prints: 0.5",
        Inches(7.2), Inches(4.7), Inches(5.8), Inches(1.1)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The big idea: self.something means \u201cthis object\u2019s something\u201d"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 7: Understanding self
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Understanding self")

    add_code_block(slide,
        "# When you write:\n"
        "sensor_a = LineSensor()\n"
        "\n"
        "# Python does this behind the scenes:\n"
        "#  1. Creates a new empty object\n"
        "#  2. Passes it to __init__ as self\n"
        "#  3. self.threshold = 0.5  sets THIS object's threshold\n"
        "#  4. Returns the object, stored in sensor_a",
        Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.5)
    )

    add_label(slide, "Rule: You write self in the definition. NEVER write self when calling.",
              Inches(0.6), Inches(5.1), Inches(12), Inches(0.5),
              font_size=Pt(17), color=DARK_BLUE)

    add_code_block(slide,
        "# Definition: self is first parameter\n"
        "def get_error(self):\n"
        "    ...\n"
        "\n"
        "# Call: no self -- Python adds it automatically\n"
        "error = sensor.get_error()",
        Inches(0.6), Inches(5.7), Inches(7.5), Inches(1.8)
    )

    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(1.5), Inches(4.8), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "self = \u201cthe object referring to itself\u201d",
                    level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Analogy:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Two students in class each say \u201cmy name is\u2026\u201d",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "self is each student\u2019s word for \u201cmy.\u201d",
                    level=0, font_size=Pt(16), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 8: Adding Methods
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Adding Methods")

    add_label(slide, "A method is a function inside a class:",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(slide,
        "class LineSensor:\n"
        "    def __init__(self):\n"
        "        self.reflectance = Reflectance.get_default_reflectance()\n"
        "        self.threshold = 0.5\n"
        "\n"
        "    def get_left(self):\n"
        "        return self.reflectance.get_left()\n"
        "\n"
        "    def get_right(self):\n"
        "        return self.reflectance.get_right()\n"
        "\n"
        "    def get_error(self):\n"
        "        return self.get_left() - self.get_right()",
        Inches(0.6), Inches(1.9), Inches(7.5), Inches(4.2)
    )

    txBox = slide.shapes.add_textbox(Inches(8.2), Inches(1.9), Inches(4.9), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Key rules for methods:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Indented inside the class (one level in)",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Always take self as the first parameter",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Use self. to access attributes (self.reflectance)",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Use self. to call other methods (self.get_left())",
                    level=0, font_size=Pt(16), color=DARK_GRAY)

    add_label(slide, "Calling methods on an object:",
              Inches(0.6), Inches(6.3), Inches(12), Inches(0.5), font_size=Pt(16))
    add_code_block(slide,
        "sensor = LineSensor()\n"
        "print(sensor.get_left())    # Read left sensor\n"
        "print(sensor.get_error())   # Compute error",
        Inches(0.6), Inches(6.8), Inches(7.5), Inches(1.3)
    )

    # ════════════════════════════════════════
    # SLIDE 9: Boolean Methods
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Boolean Methods \u2014 is_at_cross() and is_off_line()")

    add_label(slide, "Methods that return True or False:",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(slide,
        "def is_at_cross(self):\n"
        "    return self.get_left() > self.threshold and self.get_right() > self.threshold\n"
        "\n"
        "def is_off_line(self):\n"
        "    return self.get_left() < self.threshold and self.get_right() < self.threshold",
        Inches(0.6), Inches(1.9), Inches(12.0), Inches(2.1)
    )

    headers = ["Situation", "Left", "Right", "is_at_cross()", "is_off_line()"]
    rows = [
        ["On line, centered",       "high", "low",  "False", "False"],
        ["On line, left of center", "low",  "high", "False", "False"],
        ["At intersection (cross)", "high", "high", "True",  "False"],
        ["Off line entirely",       "low",  "low",  "False", "True"],
    ]
    col_widths = [Inches(3.2), Inches(1.2), Inches(1.2), Inches(2.2), Inches(2.2)]
    add_table(slide, headers, rows, Inches(0.4), Inches(4.2), Inches(10.2),
              col_widths=col_widths, row_height=Inches(0.42))

    add_code_block(slide,
        "while not sensor.is_at_cross():\n"
        "    # keep following line...\n"
        "\n"
        "if sensor.is_off_line():\n"
        '    print("Lost the line!")',
        Inches(0.4), Inches(6.0), Inches(7.5), Inches(1.5)
    )

    # ════════════════════════════════════════
    # SLIDE 10: The Complete LineSensor Class
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Complete LineSensor Class")

    add_code_block(slide,
        "from XRPLib.reflectance import Reflectance\n"
        "\n"
        "class LineSensor:\n"
        "    def __init__(self):\n"
        "        self.reflectance = Reflectance.get_default_reflectance()\n"
        "        self.threshold = 0.5\n"
        "\n"
        "    def get_left(self):\n"
        "        return self.reflectance.get_left()\n"
        "\n"
        "    def get_right(self):\n"
        "        return self.reflectance.get_right()\n"
        "\n"
        "    def get_error(self):\n"
        "        return self.get_left() - self.get_right()\n"
        "\n"
        "    def is_at_cross(self):\n"
        "        return (self.get_left() > self.threshold\n"
        "                and self.get_right() > self.threshold)\n"
        "\n"
        "    def is_off_line(self):\n"
        "        return (self.get_left() < self.threshold\n"
        "                and self.get_right() < self.threshold)",
        Inches(0.6), Inches(1.4), Inches(8.5), Inches(5.8)
    )

    txBox = slide.shapes.add_textbox(Inches(9.3), Inches(1.4), Inches(3.8), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Inventory:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "1 class", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "1 constructor (__init__)", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "2 attributes:", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "self.reflectance", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "self.threshold", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "5 methods:", level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "get_left", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "get_right", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "get_error", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "is_at_cross", level=1, font_size=Pt(15), color=DARK_GRAY)
    add_bullet_text(tf, "is_off_line", level=1, font_size=Pt(15), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 11: Before vs. After — Why Classes Win
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Before vs. After \u2014 Why Classes Win")

    add_label(slide, "Before (Lesson 7 \u2014 no class):", Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.5),
              color=RED)
    add_code_block(slide,
        "reflectance = Reflectance.get_default_reflectance()\n"
        "threshold = 0.5\n"
        "\n"
        "left     = reflectance.get_left()\n"
        "right    = reflectance.get_right()\n"
        "error    = left - right\n"
        "at_cross = left > threshold and right > threshold\n"
        "off_line = left < threshold and right < threshold",
        Inches(0.6), Inches(1.9), Inches(6.0), Inches(3.2)
    )

    add_label(slide, "After (Lesson 8 \u2014 with class):", Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
              color=GREEN)
    add_code_block(slide,
        "sensor = LineSensor()\n"
        "\n"
        "error    = sensor.get_error()\n"
        "at_cross = sensor.is_at_cross()\n"
        "off_line = sensor.is_off_line()",
        Inches(7.0), Inches(1.9), Inches(6.0), Inches(2.2)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "What improved:", level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Fewer lines of code where you use it  \u2022  Reads almost like English",
                    level=0, font_size=Pt(16), color=DARK_GRAY)
    add_bullet_text(tf, "Sensor logic is packaged and reusable  \u2022  Threshold stored in one place",
                    level=0, font_size=Pt(16), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 12: Common Mistakes & Connection to Lesson 9
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Common Mistakes \u2014 and Connection to Lesson 9")

    add_label(slide, "Watch out for these errors:",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(slide,
        "# Mistake 1: Forgetting self in method definition\n"
        "def get_error():          # WRONG -- missing self\n"
        "def get_error(self):      # CORRECT\n"
        "\n"
        "# Mistake 2: Forgetting self. when accessing attributes\n"
        "return reflectance.get_left()          # WRONG -- which reflectance?\n"
        "return self.reflectance.get_left()     # CORRECT\n"
        "\n"
        "# Mistake 3: Passing self when calling\n"
        "sensor.get_error(sensor)   # WRONG -- Python does this automatically\n"
        "sensor.get_error()         # CORRECT",
        Inches(0.6), Inches(1.9), Inches(12.0), Inches(3.8)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Coming up in Lesson 9: Object Composition",
                    level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Build a LineTrack class that uses a LineSensor inside it!",
                    level=0, font_size=Pt(17), color=DARK_GRAY)

    add_code_block(slide,
        "class LineTrack:\n"
        "    def __init__(self):\n"
        "        self.sensor    = LineSensor()  # Uses LineSensor!\n"
        "        self.drivetrain = DifferentialDrive.get_default_differential_drive()",
        Inches(0.6), Inches(6.8), Inches(10.0), Inches(1.5)
    )

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "08-introduction-to-classes.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
