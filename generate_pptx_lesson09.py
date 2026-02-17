#!/usr/bin/env python3
"""Generate PowerPoint for Module 2 Lesson 9: Object Composition — LineTrack."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(13)):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.1),
        width - Inches(0.4), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)


def add_label(slide, text, left, top, width, height,
              font_size=Pt(18), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    return txBox


def create_content_area(slide):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════
    # SLIDE 1: Title & Learning Objectives
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.0))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Module 2: Line Tracking"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson 9: Object Composition \u2014 LineTrack"
    p.font.size = Pt(38)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for obj in [
        "Understand object composition (one class using another)",
        "Build a LineTrack class that uses LineSensor",
        "Create methods for line following and turning",
        "Test classes with a simple main program",
    ]:
        add_bullet_text(tf, obj, level=0, font_size=Pt(17), color=DARK_GRAY)

    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5.5), Inches(3.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for item, duration in [
        ("Object composition concept", "10 min"),
        ("Building LineTrack", "20 min"),
        ("Testing and practice", "15 min"),
    ]:
        add_bullet_text(tf, f"{item}  ({duration})", level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 2: From Messy Code to Clean Code
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Messy Code to Clean Code")

    add_label(slide, "Lesson 7 (50+ lines, everything mixed together):",
              Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.5), font_size=Pt(17))
    add_code_block(slide,
        "while True:\n"
        "    left  = reflectance.get_left()\n"
        "    right = reflectance.get_right()\n"
        "    if left > 0.5 and right > 0.5:\n"
        "        drivetrain.stop()\n"
        "        drivetrain.turn(180)\n"
        "        time.sleep(0.3)\n"
        "    else:\n"
        "        error = left - right\n"
        "        # ... more code ...",
        Inches(0.6), Inches(1.9), Inches(6.0), Inches(3.5)
    )

    add_label(slide, "What we want (3 lines!):",
              Inches(7.2), Inches(1.3), Inches(5.8), Inches(0.5), font_size=Pt(17))
    add_code_block(slide,
        "tracker = LineTrack()\n"
        "tracker.track_until_cross()\n"
        "tracker.turn_right()",
        Inches(7.2), Inches(1.9), Inches(5.8), Inches(1.6)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Classes hide complexity behind simple names."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = BODY_FONT

    # ════════════════════════════════════════
    # SLIDE 3: What Is Object Composition?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "What Is Object Composition?")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Composition: One class contains and uses another class\u2019s object.",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Real-world analogy:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "A car HAS AN engine and HAS wheels",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "The car doesn\u2019t build the engine \u2014 it uses one",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "The car delegates \u201cgo fast\u201d to the engine",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "In our code:", level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "LineTrack HAS A LineSensor and HAS A DifferentialDrive",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "LineTrack delegates \u201ccheck sensor\u201d to LineSensor",
                    level=1, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 4: LineTrack.__init__()
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "LineTrack.__init__()")

    add_code_block(slide,
        "class LineTrack:\n"
        "    def __init__(self):\n"
        "        self.sensor     = LineSensor()       # Create a LineSensor\n"
        "        self.drivetrain = DifferentialDrive.get_default_differential_drive()\n"
        "        self.base_effort = 0.4\n"
        "        self.Kp          = 0.5",
        Inches(0.6), Inches(1.5), Inches(12.0), Inches(2.5)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "self.sensor = LineSensor()",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Creates a LineSensor object and stores it inside LineTrack.",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Now LineTrack can call:",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "self.sensor.get_error()",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "self.sensor.is_at_cross()",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "self.sensor.is_off_line()",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 5: track_until_cross()
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "track_until_cross()")

    add_code_block(slide,
        "def track_until_cross(self):\n"
        "    while not self.sensor.is_at_cross():\n"
        "        error = self.sensor.get_error()\n"
        "        left  = self.base_effort - error * self.Kp\n"
        "        right = self.base_effort + error * self.Kp\n"
        "        self.drivetrain.set_effort(left, right)\n"
        "    self.drivetrain.stop()",
        Inches(0.6), Inches(1.5), Inches(10.0), Inches(3.0)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "while not self.sensor.is_at_cross():",
                    level=0, font_size=Pt(20), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Keep following while NOT at a cross",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "When cross detected \u2192 loop exits \u2192 stop",
                    level=1, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Uses LineSensor methods for ALL sensor work.",
                    level=0, font_size=Pt(18), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 6: turn_right() and turn_left()
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "turn_right() and turn_left()")

    add_code_block(slide,
        "def turn_right(self):\n"
        "    # Drive forward to clear intersection\n"
        "    self.drivetrain.set_effort(self.base_effort, self.base_effort)\n"
        "    time.sleep(0.3)\n"
        "    # Turn right\n"
        "    self.drivetrain.set_effort(0.3, -0.3)\n"
        "    time.sleep(0.3)\n"
        "    # Keep turning until line is found\n"
        "    while self.sensor.is_off_line():\n"
        "        pass\n"
        "    self.drivetrain.stop()",
        Inches(0.6), Inches(1.5), Inches(9.5), Inches(4.0)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Three phases: Clear intersection \u2192 start turning \u2192 find line again",
                    level=0, font_size=Pt(18), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "turn_left() is the same but with opposite motor directions.",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 7: Using LineTrack in a Main Program
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Using LineTrack in a Main Program")

    add_code_block(slide,
        "board   = Board.get_default_board()\n"
        "tracker = LineTrack()\n"
        "\n"
        "board.wait_for_button()\n"
        "\n"
        "tracker.track_until_cross()\n"
        'print("Cross found!")\n'
        "tracker.turn_right()\n"
        'print("Turned right!")\n'
        "tracker.track_until_cross()\n"
        'print("Done!")',
        Inches(0.6), Inches(1.5), Inches(8.5), Inches(4.5)
    )

    txBox = slide.shapes.add_textbox(Inches(9.2), Inches(1.5), Inches(3.9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "The main program is clean and readable.",
                    level=0, font_size=Pt(17), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "All complex sensor/motor code is hidden inside the classes.",
                    level=0, font_size=Pt(17), color=DARK_GRAY)

    # ════════════════════════════════════════
    # SLIDE 8: The Building Block Pattern
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Building Block Pattern")

    add_code_block(slide,
        "LineSensor          DifferentialDrive\n"
        "    \\                      /\n"
        "     \\                    /\n"
        "      \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "             \u2193\n"
        "         LineTrack\n"
        "             \u2193\n"
        "        Main Program",
        Inches(2.0), Inches(1.5), Inches(9.0), Inches(4.0),
        font_size=Pt(20)
    )

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Each layer builds on the layer below.",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "This is how professional software is built!",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ════════════════════════════════════════
    # SLIDE 9: Your Turn!
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Your Turn!")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "Exercise 1:", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Build LineTrack and test each method",
                    level=0, font_size=Pt(18), color=DARK_GRAY)
    add_bullet_text(tf, "track_until_cross() \u2014 does it follow and stop?",
                    level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "turn_right() \u2014 does it turn and find the line?",
                    level=1, font_size=Pt(17), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Exercise 2 (Challenge):", level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Chain multiple maneuvers",
                    level=0, font_size=Pt(18), color=DARK_GRAY)

    add_code_block(slide,
        "for i in range(4):\n"
        "    tracker.track_until_cross()\n"
        "    tracker.turn_right()",
        Inches(7.2), Inches(2.5), Inches(5.8), Inches(1.8)
    )

    # ════════════════════════════════════════
    # SLIDE 10: Connection to Final Project
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Connection to Final Project")

    tf = create_content_area(slide)
    p = tf.paragraphs[0]
    p.text = ""
    add_bullet_text(tf, "You now have two reusable classes:",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "LineSensor \u2014 reads and interprets sensors",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "LineTrack \u2014 follows lines and handles intersections",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "", level=0)
    add_bullet_text(tf, "Next (Lesson 10 \u2014 Final Project):",
                    level=0, font_size=Pt(22), bold=True, color=DARK_BLUE)
    add_bullet_text(tf, "Follow circle \u2192 detect cross \u2192 reverse \u2192 repeat 4 times",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "Uses BOTH classes together",
                    level=0, font_size=Pt(20), color=DARK_GRAY)
    add_bullet_text(tf, "The capstone of Module 2!",
                    level=0, font_size=Pt(20), bold=True, color=ACCENT_BLUE)

    # ── Save ──
    output_path = os.path.join(
        os.path.dirname(__file__),
        "module-02-line-tracking", "slides", "09-object-composition.pptx"
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
