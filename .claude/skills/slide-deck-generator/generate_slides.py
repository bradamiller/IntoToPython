#!/usr/bin/env python3
"""Generate PowerPoint and HTML slide decks from markdown outline files.

Usage:
    python generate_slides.py <markdown_file> [<output_file>] [--module "Module Name"] [--html]
    python generate_slides.py --batch <directory> [--module "Module Name"] [--html]

Examples:
    python generate_slides.py module-02-line-tracking/slides/01-the-reflectance-sensor-outline.md
    python generate_slides.py --batch module-01-driving/slides --html
"""

import re
import sys
import os
import glob as globmod
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
GREEN = RGBColor(0x17, 0x7B, 0x41)
RED = RGBColor(0xC0, 0x00, 0x00)

# ── Fonts ──
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
CODE_FONT = "Courier New"

# ── Sizes ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MIN_CODE_FONT_SIZE = Pt(14)  # Improved: minimum readable code font size


def parse_markdown_slides(md_text):
    """Parse a slide outline markdown into structured slide data."""
    slide_pattern = re.compile(r'^## Slide \d+:\s*(.+)', re.MULTILINE)
    parts = slide_pattern.split(md_text)

    slides = []
    lesson_title = ""
    preamble = parts[0].strip()
    m = re.match(r'^#\s+(.+)', preamble)
    if m:
        lesson_title = m.group(1).strip()

    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        body = parts[i + 1] if i + 1 < len(parts) else ""
        body = re.sub(r'\n---\s*$', '', body.strip())

        blocks = parse_slide_body(body)
        slides.append({
            'number': (i + 1) // 2,
            'title': title,
            'content_blocks': blocks,
        })

    return lesson_title, slides


def parse_slide_body(body):
    """Parse the body of a slide into content blocks."""
    blocks = []
    lines = body.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]

        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append({'type': 'code', 'code': '\n'.join(code_lines)})
            continue

        if '|' in line and i + 1 < len(lines) and '---' in lines[i + 1]:
            headers = [c.strip() for c in line.strip().strip('|').split('|')]
            i += 2
            rows = []
            while i < len(lines) and '|' in lines[i]:
                row = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                rows.append(row)
                i += 1
            blocks.append({'type': 'table', 'headers': headers, 'rows': rows})
            continue

        if line.strip() == '':
            blocks.append({'type': 'blank'})
            i += 1
            continue

        stripped = line.strip()
        is_bullet = False
        bullet_level = 0
        bullet_text = stripped

        if re.match(r'^[-*]\s', stripped):
            is_bullet = True
            bullet_text = re.sub(r'^[-*]\s+', '', stripped)
        elif re.match(r'^\d+\.\s', stripped):
            is_bullet = True
            bullet_text = re.sub(r'^\d+\.\s+', '', stripped)
        elif re.match(r'^\s+[-*]\s', line):
            is_bullet = True
            bullet_level = 1
            bullet_text = re.sub(r'^[-*]\s+', '', stripped)
        elif re.match(r'^\s+\d+\.\s', line):
            is_bullet = True
            bullet_level = 1
            bullet_text = re.sub(r'^\d+\.\s+', '', stripped)

        clean_text = clean_markdown(bullet_text)
        is_bold = stripped.startswith('**') and ':**' in stripped[:50]
        is_header_line = stripped.startswith('**') and stripped.endswith('**') and not is_bullet

        if is_bullet:
            blocks.append({
                'type': 'bullet',
                'text': clean_text,
                'level': bullet_level,
                'bold': is_bold,
            })
        else:
            blocks.append({
                'type': 'text',
                'text': clean_text,
                'bold': is_bold or is_header_line,
                'is_header': is_bold or is_header_line,
            })

        i += 1

    return blocks


def clean_markdown(text):
    """Remove markdown formatting characters."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = text.replace('→', '\u2192')
    text = text.replace('—', '\u2014')
    text = text.replace('–', '\u2013')
    return text


# ═══════════════════════════════════════════
# PPTX GENERATION HELPERS
# ═══════════════════════════════════════════

def add_title_bar(slide, title_text):
    """Add a colored title bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT


def add_bullet_text(text_frame, text, level=0, font_size=Pt(18), bold=False, color=BLACK):
    """Add a bullet point to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = BODY_FONT
    p.font.bold = bold
    p.space_after = Pt(4)
    return p


def add_code_block(slide, code_text, left, top, width, height):
    """Add a code block with gray background and readable fonts."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = MEDIUM_GRAY
    bg.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.1),
        width - Inches(0.4), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code_text.strip().split('\n')
    # Calculate optimal font size for readability
    max_lines = len(lines)
    available_height = height - Inches(0.3)
    line_height = available_height / max_lines

    # Ensure minimum readable font size
    font_size = max(MIN_CODE_FONT_SIZE, Pt(min(18, line_height * 72)))  # Convert inches to points

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.4)):
    """Add a formatted table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, row_height * num_rows
    )
    table = table_shape.table

    col_width = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = BODY_FONT
            p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════

def build_title_slide(prs, slide_data, lesson_title, module_name):
    """Build Slide 1: Title & Learning Objectives."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(3.0)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = module_name
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.font.name = TITLE_FONT

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    m = re.search(r':\s*(.+)', lesson_title)
    if m:
        display_title = m.group(1).strip()
    else:
        display_title = lesson_title.replace("Slide Outline: ", "").replace("Lesson ", "")
    p.text = display_title
    p.font.size = Pt(42)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = TITLE_FONT

    blocks = slide_data['content_blocks']
    objectives = []
    agenda_items = []
    current_section = None

    for block in blocks:
        text = block.get('text', '').lower()
        if 'learning objective' in text:
            current_section = 'objectives'
            continue
        elif 'agenda' in text:
            current_section = 'agenda'
            continue
        elif block.get('is_header') or (block['type'] == 'text' and block.get('bold')):
            if current_section not in ('objectives', 'agenda'):
                current_section = None

        if block['type'] in ('bullet', 'text') and block.get('text', '').strip():
            if current_section == 'objectives':
                objectives.append(block['text'])
            elif current_section == 'agenda':
                agenda_items.append(block['text'])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = TITLE_FONT

    for obj in objectives:
        add_bullet_text(tf, obj, level=0, font_size=Pt(16), color=DARK_GRAY)

    if agenda_items:
        txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(5), Inches(3.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = TITLE_FONT

        for item in agenda_items:
            add_bullet_text(tf, item, level=0, font_size=Pt(16), color=DARK_GRAY)


def build_content_slide(prs, slide_data):
    """Build a generic content slide from parsed blocks."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_bar(slide, slide_data['title'])

    blocks = slide_data['content_blocks']

    code_blocks = [b for b in blocks if b['type'] == 'code']
    tables = [b for b in blocks if b['type'] == 'table']
    text_blocks = [b for b in blocks if b['type'] in ('text', 'bullet')]

    if len(code_blocks) >= 2 and len(text_blocks) <= 8:
        build_dual_code_layout(slide, blocks)
    elif len(code_blocks) == 1 and len(text_blocks) > 4:
        build_code_and_text_layout(slide, blocks)
    elif tables:
        build_table_layout(slide, blocks)
    else:
        build_sequential_layout(slide, blocks)


def build_sequential_layout(slide, blocks):
    """Default layout: text and code blocks stacked vertically with overflow protection."""
    y_pos = Inches(1.5)
    x_left = Inches(0.6)
    content_width = Inches(12)
    max_y = Inches(6.5)  # Prevent overflow at bottom

    text_tf = None
    for block in blocks:
        if block['type'] == 'blank':
            if text_tf:
                add_bullet_text(text_tf, '', level=0, font_size=Pt(4))
            continue

        if block['type'] == 'code':
            if text_tf:
                text_tf = None

            code_lines = block['code'].strip().split('\n')
            code_height = max(Inches(0.8), Inches(0.25 * len(code_lines) + 0.3))
            code_height = min(code_height, max_y - y_pos - Inches(0.5))  # Prevent overflow

            if y_pos + code_height > max_y:
                break  # Skip if would overflow

            add_code_block(slide, block['code'], x_left, y_pos, content_width, code_height)
            y_pos += code_height + Inches(0.15)
            continue

        if block['type'] == 'table':
            if text_tf:
                text_tf = None
            headers = block['headers']
            rows = block['rows']
            table_height = Inches(0.4) * (len(rows) + 1)

            if y_pos + table_height > max_y:
                break  # Skip if would overflow

            add_table(slide, headers, rows, Inches(1.0), y_pos, Inches(11), Inches(0.4))
            y_pos += table_height + Inches(0.15)
            continue

        # Text or bullet
        if text_tf is None:
            txBox = slide.shapes.add_textbox(x_left, y_pos, content_width, Inches(0.4))
            text_tf = txBox.text_frame
            text_tf.word_wrap = True
            p = text_tf.paragraphs[0]
            text = block.get('text', '')
            is_header = block.get('is_header', False) or block.get('bold', False)
            p.text = text
            p.font.size = Pt(20) if is_header else Pt(17)
            p.font.bold = is_header
            p.font.color.rgb = DARK_BLUE if is_header else BLACK
            p.font.name = BODY_FONT
            p.level = block.get('level', 0) if block['type'] == 'bullet' else 0
            p.space_after = Pt(4)
            y_pos += Inches(0.3)
        else:
            if y_pos > max_y - Inches(0.5):
                break  # Prevent overflow
            text = block.get('text', '')
            is_header = block.get('is_header', False) or block.get('bold', False)
            level = block.get('level', 0) if block['type'] == 'bullet' else 0
            add_bullet_text(
                text_tf, text,
                level=level,
                font_size=Pt(20) if is_header else Pt(17),
                bold=is_header,
                color=DARK_BLUE if is_header else BLACK,
            )
            y_pos += Inches(0.2)


def build_dual_code_layout(slide, blocks):
    """Layout with two code blocks side by side with overflow protection."""
    sections = []
    current_text = []
    code_blocks = []

    for block in blocks:
        if block['type'] == 'code':
            if current_text:
                sections.append(('text', current_text))
                current_text = []
            code_blocks.append(block)
        elif block['type'] == 'table':
            if current_text:
                sections.append(('text', current_text))
                current_text = []
            sections.append(('table', block))
        else:
            current_text.append(block)
    if current_text:
        sections.append(('text', current_text))

    y_pos = Inches(1.5)
    max_y = Inches(6.5)

    for section_type, section_data in sections:
        if section_type == 'text':
            txBox = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(12), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            first = True
            count = 0
            for block in section_data:
                if block['type'] == 'blank':
                    continue
                text = block.get('text', '')
                if not text.strip():
                    continue
                is_header = block.get('is_header', False) or block.get('bold', False)
                if first:
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(18) if is_header else Pt(16)
                    p.font.bold = is_header
                    p.font.color.rgb = DARK_BLUE if is_header else BLACK
                    p.font.name = BODY_FONT
                    first = False
                else:
                    add_bullet_text(
                        tf, text,
                        level=block.get('level', 0) if block['type'] == 'bullet' else 0,
                        font_size=Pt(18) if is_header else Pt(16),
                        bold=is_header,
                        color=DARK_BLUE if is_header else BLACK,
                    )
                count += 1
            y_pos += Inches(0.3) * max(1, count)

    if len(code_blocks) >= 2:
        for idx in range(0, len(code_blocks), 2):
            cb1 = code_blocks[idx]
            lines1 = cb1['code'].strip().split('\n')
            h1 = max(Inches(0.8), Inches(0.25 * len(lines1) + 0.3))

            if idx + 1 < len(code_blocks):
                cb2 = code_blocks[idx + 1]
                lines2 = cb2['code'].strip().split('\n')
                h2 = max(Inches(0.8), Inches(0.25 * len(lines2) + 0.3))
                h = min(max(h1, h2), max_y - y_pos - Inches(0.5))
                if y_pos + h > max_y:
                    break
                add_code_block(slide, cb1['code'], Inches(0.6), y_pos, Inches(5.8), h)
                add_code_block(slide, cb2['code'], Inches(6.8), y_pos, Inches(5.8), h)
            else:
                h = min(h1, max_y - y_pos - Inches(0.5))
                if y_pos + h > max_y:
                    break
                add_code_block(slide, cb1['code'], Inches(0.6), y_pos, Inches(12), h)
            y_pos += h + Inches(0.15)


def build_code_and_text_layout(slide, blocks):
    """Layout with code block and text side by side with overflow protection."""
    code_block = None
    text_before = []
    text_after = []
    found_code = False

    for block in blocks:
        if block['type'] == 'code' and code_block is None:
            code_block = block
            found_code = True
        elif block['type'] == 'table':
            if not found_code:
                text_before.append(block)
            else:
                text_after.append(block)
        elif not found_code:
            text_before.append(block)
        else:
            text_after.append(block)

    y_pos = Inches(1.5)
    max_y = Inches(6.5)

    if text_before:
        txBox = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(12), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        count = 0
        for block in text_before:
            if block['type'] == 'blank':
                continue
            text = block.get('text', '')
            if not text.strip():
                continue
            is_header = block.get('is_header', False) or block.get('bold', False)
            if first:
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(18) if is_header else Pt(16)
                p.font.bold = is_header
                p.font.color.rgb = DARK_BLUE if is_header else BLACK
                p.font.name = BODY_FONT
                first = False
            else:
                add_bullet_text(
                    tf, text,
                    level=block.get('level', 0) if block['type'] == 'bullet' else 0,
                    font_size=Pt(18) if is_header else Pt(16),
                    bold=is_header,
                    color=DARK_BLUE if is_header else BLACK,
                )
            count += 1
        y_pos += Inches(0.3) * max(1, count)

    if code_block and y_pos < max_y:
        code_lines = code_block['code'].strip().split('\n')
        code_h = min(max(Inches(0.8), Inches(0.25 * len(code_lines) + 0.3)), max_y - y_pos)
        add_code_block(slide, code_block['code'], Inches(0.6), y_pos, Inches(6.2), code_h)

        if text_after:
            txBox = slide.shapes.add_textbox(Inches(7.2), y_pos, Inches(5.5), code_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            first = True
            for block in text_after:
                if block['type'] == 'blank':
                    continue
                text = block.get('text', '')
                if not text.strip():
                    continue
                is_header = block.get('is_header', False) or block.get('bold', False)
                if first:
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(18) if is_header else Pt(16)
                    p.font.bold = is_header
                    p.font.color.rgb = DARK_BLUE if is_header else BLACK
                    p.font.name = BODY_FONT
                    first = False
                else:
                    add_bullet_text(
                        tf, text,
                        level=block.get('level', 0) if block['type'] == 'bullet' else 0,
                        font_size=Pt(18) if is_header else Pt(16),
                        bold=is_header,
                        color=DARK_BLUE if is_header else BLACK,
                    )


def build_table_layout(slide, blocks):
    """Layout that handles tables properly with overflow protection."""
    y_pos = Inches(1.5)
    max_y = Inches(6.5)
    text_tf = None

    for block in blocks:
        if block['type'] == 'blank':
            if text_tf:
                add_bullet_text(text_tf, '', level=0, font_size=Pt(4))
            continue

        if block['type'] == 'table':
            text_tf = None
            headers = block['headers']
            rows = block['rows']
            table_height = Inches(0.4) * (len(rows) + 1)

            if y_pos + table_height > max_y:
                break

            add_table(slide, headers, rows, Inches(1.0), y_pos, Inches(11))
            y_pos += table_height + Inches(0.2)
            continue

        if block['type'] == 'code':
            text_tf = None
            code_lines = block['code'].strip().split('\n')
            code_h = min(max(Inches(0.8), Inches(0.25 * len(code_lines) + 0.3)), max_y - y_pos - Inches(0.5))

            if y_pos + code_h > max_y:
                break

            add_code_block(slide, block['code'], Inches(0.6), y_pos, Inches(12), code_h)
            y_pos += code_h + Inches(0.15)
            continue

        text = block.get('text', '')
        is_header = block.get('is_header', False) or block.get('bold', False)
        level = block.get('level', 0) if block['type'] == 'bullet' else 0

        if text_tf is None:
            txBox = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(12), Inches(0.4))
            text_tf = txBox.text_frame
            text_tf.word_wrap = True
            p = text_tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(18) if is_header else Pt(16)
            p.font.bold = is_header
            p.font.color.rgb = DARK_BLUE if is_header else BLACK
            p.font.name = BODY_FONT
            p.level = level
            p.space_after = Pt(4)
        else:
            if y_pos > max_y - Inches(0.5):
                break
            add_bullet_text(
                text_tf, text,
                level=level,
                font_size=Pt(18) if is_header else Pt(16),
                bold=is_header,
                color=DARK_BLUE if is_header else BLACK,
            )


# ═══════════════════════════════════════════
# HTML GENERATION
# ═══════════════════════════════════════════

def generate_html_slides(md_path, output_path=None, module_name=None):
    """Generate HTML slides with navigation from markdown outline."""
    with open(md_path, 'r') as f:
        md_text = f.read()

    lesson_title, slides = parse_markdown_slides(md_text)

    if module_name is None:
        if 'module-01' in md_path:
            module_name = "Module 1: Driving"
        elif 'module-02' in md_path:
            module_name = "Module 2: Line Tracking"
        elif 'module-03' in md_path:
            module_name = "Module 3: Rangefinder & Navigation"
        elif 'module-04' in md_path:
            module_name = "Module 4: Data & Algorithms"
        elif 'module-05' in md_path:
            module_name = "Module 5: Capstone"
        else:
            module_name = "XRP Python Programming"

    if output_path is None:
        output_path = md_path.replace('-outline.md', '.html')

    # Generate slide HTML
    slide_htmls = []
    for slide_data in slides:
        slide_html = generate_slide_html(slide_data, lesson_title, module_name, len(slides))
        slide_htmls.append(slide_html)

    # Generate navigation
    nav_items = []
    for i, slide_data in enumerate(slides):
        nav_items.append(f'<a href="#slide-{i+1}" class="nav-item" onclick="showSlide({i})">{slide_data["title"]}</a>')

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{lesson_title} - {module_name}</title>
    <style>
        {get_css_styles()}
    </style>
</head>
<body>
    <div class="nav-sidebar">
        <div class="nav-header">
            <h2>{module_name}</h2>
            <h3>{lesson_title}</h3>
        </div>
        <nav class="slide-nav">
            {chr(10).join(nav_items)}
        </nav>
    </div>
    <div class="slides-container">
        {chr(10).join(slide_htmls)}
    </div>
    <script>
        {get_js_script(len(slides))}
    </script>
</body>
</html>"""

    with open(output_path, 'w') as f:
        f.write(html_content)

    print(f"  Created: {output_path} ({len(slides)} slides)")
    return output_path


def generate_slide_html(slide_data, lesson_title, module_name, total_slides):
    """Generate HTML for a single slide."""
    slide_num = slide_data['number']
    title = slide_data['title']
    blocks = slide_data['content_blocks']

    content_html = ""

    for block in blocks:
        if block['type'] == 'text':
            css_class = "header-text" if block.get('is_header') else "body-text"
            bold_class = " bold" if block.get('bold') else ""
            content_html += f'<p class="{css_class}{bold_class}">{block["text"]}</p>\n'
        elif block['type'] == 'bullet':
            level_class = f" level-{block.get('level', 0)}"
            bold_class = " bold" if block.get('bold') else ""
            content_html += f'<p class="bullet{level_class}{bold_class}">{block["text"]}</p>\n'
        elif block['type'] == 'code':
            content_html += f'<pre class="code-block"><code>{block["code"]}</code></pre>\n'
        elif block['type'] == 'table':
            table_html = '<table class="content-table">\n<thead>\n<tr>\n'
            for header in block['headers']:
                table_html += f'<th>{header}</th>\n'
            table_html += '</tr>\n</thead>\n<tbody>\n'
            for row in block['rows']:
                table_html += '<tr>\n'
                for cell in row:
                    table_html += f'<td>{cell}</td>\n'
                table_html += '</tr>\n'
            table_html += '</tbody>\n</table>\n'
            content_html += table_html
        elif block['type'] == 'blank':
            content_html += '<div class="spacer"></div>\n'

    prev_disabled = "disabled" if slide_num == 1 else ""
    next_disabled = "disabled" if slide_num == total_slides else ""

    return f'''<div id="slide-{slide_num}" class="slide">
    <div class="title-bar">
        <h1>{title}</h1>
    </div>
    <div class="content">
        {content_html}
    </div>
    <div class="slide-navigation">
        <button onclick="prevSlide()" class="nav-btn prev-btn" {prev_disabled}>Previous</button>
        <span class="slide-counter">{slide_num} / {total_slides}</span>
        <button onclick="nextSlide()" class="nav-btn next-btn" {next_disabled}>Next</button>
    </div>
</div>'''


def get_css_styles():
    """Return CSS styles for HTML slides."""
    return """
* {
    margin: 0;
    padding: 0;
    box-sizing: border-box;
}

body {
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    background-color: #f5f5f5;
    overflow: hidden;
}

.nav-sidebar {
    position: fixed;
    left: 0;
    top: 0;
    width: 250px;
    height: 100vh;
    background-color: #1B3A5C;
    color: white;
    padding: 20px;
    overflow-y: auto;
    z-index: 1000;
}

.nav-header h2 {
    font-size: 18px;
    margin-bottom: 10px;
    color: #A0C4E8;
}

.nav-header h3 {
    font-size: 14px;
    margin-bottom: 20px;
    color: #FFFFFF;
}

.slide-nav {
    display: flex;
    flex-direction: column;
}

.nav-item {
    color: #FFFFFF;
    text-decoration: none;
    padding: 8px 12px;
    margin-bottom: 4px;
    border-radius: 4px;
    font-size: 13px;
    line-height: 1.4;
    transition: background-color 0.2s;
}

.nav-item:hover {
    background-color: #2E75B6;
}

.slides-container {
    margin-left: 250px;
    height: 100vh;
    position: relative;
}

.slide {
    position: absolute;
    top: 0;
    left: 0;
    width: calc(100vw - 250px);
    height: 100vh;
    background-color: white;
    display: none;
    flex-direction: column;
}

.slide.active {
    display: flex;
}

.title-bar {
    background-color: #1B3A5C;
    color: white;
    padding: 20px 40px;
    height: 80px;
    display: flex;
    align-items: center;
}

.title-bar h1 {
    font-size: 28px;
    font-weight: bold;
}

.content {
    flex: 1;
    padding: 20px;
    padding-bottom: 80px;  /* Make room for navigation */
    overflow-y: auto;
    max-width: 1200px;
}

.header-text {
    font-size: 24px;
    font-weight: bold;
    color: #1B3A5C;
    margin-bottom: 16px;
}

.body-text {
    font-size: 18px;
    margin-bottom: 12px;
    line-height: 1.5;
}

.bold {
    font-weight: bold;
}

.bullet {
    font-size: 18px;
    margin-bottom: 8px;
    margin-left: 20px;
    line-height: 1.4;
}

.bullet.level-1 {
    margin-left: 40px;
}

.code-block {
    background-color: #F0F0F0;
    border: 1px solid #999;
    border-radius: 8px;
    padding: 16px;
    margin: 16px 0;
    font-family: 'Courier New', monospace;
    font-size: 14px;
    line-height: 1.4;
    color: #333;
    overflow-x: auto;
}

.content-table {
    width: 100%;
    border-collapse: collapse;
    margin: 16px 0;
    font-size: 14px;
}

.content-table th {
    background-color: #2E75B6;
    color: white;
    padding: 8px 12px;
    text-align: center;
    font-weight: bold;
}

.content-table td {
    padding: 6px 12px;
    text-align: center;
    border-bottom: 1px solid #ddd;
}

.content-table tr:nth-child(even) {
    background-color: #E8F0F8;
}

.spacer {
    height: 16px;
}

.slide-navigation {
    position: absolute;
    bottom: 20px;
    left: 50%;
    transform: translateX(-50%);
    display: flex;
    align-items: center;
    gap: 20px;
    background-color: rgba(255, 255, 255, 0.9);
    padding: 10px 20px;
    border-radius: 25px;
    box-shadow: 0 2px 10px rgba(0, 0, 0, 0.1);
}

.nav-btn {
    padding: 8px 16px;
    border: none;
    border-radius: 20px;
    background-color: #2E75B6;
    color: white;
    font-size: 14px;
    cursor: pointer;
    transition: background-color 0.2s;
}

.nav-btn:hover:not(:disabled) {
    background-color: #1B3A5C;
}

.nav-btn:disabled {
    background-color: #ccc;
    cursor: not-allowed;
}

.slide-counter {
    font-size: 14px;
    font-weight: bold;
    color: #1B3A5C;
    min-width: 60px;
    text-align: center;
}
"""
    """Return JavaScript for slide navigation."""
    return f"""
let currentSlide = 0;
const slides = document.querySelectorAll('.slide');

function showSlide(index) {{
    slides.forEach((slide, i) => {{
        if (i === index) {{
            slide.classList.add('active');
        }} else {{
            slide.classList.remove('active');
        }}
    }});
    currentSlide = index;

    // Update navigation highlighting
    document.querySelectorAll('.nav-item').forEach((item, i) => {{
        if (i === index) {{
            item.classList.add('active');
        }} else {{
            item.classList.remove('active');
        }}
    }});

    // Update button states
    const prevBtn = document.querySelector('.slide.active .prev-btn');
    const nextBtn = document.querySelector('.slide.active .next-btn');
    if (prevBtn) prevBtn.disabled = index === 0;
    if (nextBtn) nextBtn.disabled = index === slides.length - 1;
}}

function nextSlide() {{
    if (currentSlide < {num_slides - 1}) {{
        showSlide(currentSlide + 1);
    }}
}}

function prevSlide() {{
    if (currentSlide > 0) {{
        showSlide(currentSlide - 1);
    }}
}}

// Keyboard navigation
document.addEventListener('keydown', (e) => {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{
        e.preventDefault();
        nextSlide();
    }} else if (e.key === 'ArrowLeft') {{
        e.preventDefault();
        prevSlide();
    }}
}});

// Initialize first slide
showSlide(0);
"""


def get_js_script(num_slides):
    """Return JavaScript for slide navigation."""
    return f"""
let currentSlide = 0;
const slides = document.querySelectorAll('.slide');

function showSlide(index) {{
    slides.forEach((slide, i) => {{
        if (i === index) {{
            slide.classList.add('active');
        }} else {{
            slide.classList.remove('active');
        }}
    }});
    currentSlide = index;

    // Update navigation highlighting
    document.querySelectorAll('.nav-item').forEach((item, i) => {{
        if (i === index) {{
            item.classList.add('active');
        }} else {{
            item.classList.remove('active');
        }}
    }});

    // Update button states
    const prevBtn = document.querySelector('.slide.active .prev-btn');
    const nextBtn = document.querySelector('.slide.active .next-btn');
    if (prevBtn) prevBtn.disabled = index === 0;
    if (nextBtn) nextBtn.disabled = index === slides.length - 1;
}}

function nextSlide() {{
    if (currentSlide < {num_slides - 1}) {{
        showSlide(currentSlide + 1);
    }}
}}

function prevSlide() {{
    if (currentSlide > 0) {{
        showSlide(currentSlide - 1);
    }}
}}

// Keyboard navigation
document.addEventListener('keydown', (e) => {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{
        e.preventDefault();
        nextSlide();
    }} else if (e.key === 'ArrowLeft') {{
        e.preventDefault();
        prevSlide();
    }}
}});

// Initialize first slide
showSlide(0);
"""


# ═══════════════════════════════════════════
# MAIN BUILD FUNCTIONS
# ═══════════════════════════════════════════

def build_presentation(md_path, output_path=None, module_name=None, html=False):
    """Build a PPTX or HTML presentation from a markdown slide outline file."""
    if html:
        return generate_html_slides(md_path, output_path, module_name)
    else:
        with open(md_path, 'r') as f:
            md_text = f.read()

        lesson_title, slides = parse_markdown_slides(md_text)

        if module_name is None:
            if 'module-01' in md_path:
                module_name = "Module 1: Driving"
            elif 'module-02' in md_path:
                module_name = "Module 2: Line Tracking"
            elif 'module-03' in md_path:
                module_name = "Module 3: Rangefinder & Navigation"
            elif 'module-04' in md_path:
                module_name = "Module 4: Data & Algorithms"
            elif 'module-05' in md_path:
                module_name = "Module 5: Capstone"
            else:
                module_name = "XRP Python Programming"

        if output_path is None:
            output_path = md_path.replace('-outline.md', '.pptx')

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        for slide_data in slides:
            if slide_data['number'] == 1:
                build_title_slide(prs, slide_data, lesson_title, module_name)
            else:
                build_content_slide(prs, slide_data)

        prs.save(output_path)
        print(f"  Created: {output_path} ({len(prs.slides)} slides)")
        return output_path


def batch_generate(directory, module_name=None, html=False):
    """Generate PPTX or HTML files for all outline markdown files in a directory."""
    pattern = os.path.join(directory, '*-outline.md')
    files = sorted(globmod.glob(pattern))

    files = [f for f in files if 'structure' not in os.path.basename(f)]

    if not files:
        print(f"No outline files found in {directory}")
        return

    print(f"Found {len(files)} slide outlines in {directory}")
    print()

    created = []
    for md_path in files:
        name = os.path.basename(md_path)
        print(f"Processing: {name}")
        out = build_presentation(md_path, module_name=module_name, html=html)
        created.append(out)

    print(f"\nDone! Created {len(created)} presentations.")
    return created


if __name__ == "__main__":
    args = sys.argv[1:]

    if not args:
        print(__doc__)
        sys.exit(1)

    module_name = None
    html = False

    # Parse flags
    if '--module' in args:
        idx = args.index('--module')
        if idx + 1 < len(args):
            module_name = args[idx + 1]
            args = args[:idx] + args[idx + 2:]

    if '--html' in args:
        html = True
        args.remove('--html')

    if args[0] == '--batch':
        if len(args) < 2:
            print("Usage: python generate_slides.py --batch <directory> [--module 'Module Name'] [--html]")
            sys.exit(1)
        batch_generate(args[1], module_name=module_name, html=html)
    else:
        md_path = args[0]
        output_path = args[1] if len(args) > 1 else None
        build_presentation(md_path, output_path, module_name=module_name, html=html)